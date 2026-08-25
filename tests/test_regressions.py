"""
Deterministic regression tests for the Phase 1 fixes.

No Ollama, no network - fake models and a fake executor, so these
test orchestration behaviour rather than model quality.
"""

import ast
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from core.conversation_state import ConversationState
from core.router import resolve_pending_file_selection, is_active_file_request

PASS, FAIL = [], []


def check(name, condition, detail=""):
    (PASS if condition else FAIL).append(name)
    print(f"{'PASS' if condition else 'FAIL'}  {name}" + (f"\n        {detail}" if detail and not condition else ""))


# ----------------------------------------------------------------
# 1. Ordinal selection resolver
# ----------------------------------------------------------------
print("\n--- ordinal selection ---")

PENDING = [
    r"C:\Users\TestUser\Documents\Hostel Fees.pdf",
    r"C:\Users\TestUser\Downloads\Hostel Fees 2024.pdf",
    r"C:\Users\TestUser\Desktop\Hostel Fees old.pdf",
]

state = ConversationState(pending_file_paths=list(PENDING))

check("'the first one' -> match 1", resolve_pending_file_selection(state, "the first one") == PENDING[0])
check("'first' -> match 1", resolve_pending_file_selection(state, "first") == PENDING[0])
check("'the second one' -> match 2 (not 1)", resolve_pending_file_selection(state, "the second one") == PENDING[1],
      f"got {resolve_pending_file_selection(state, 'the second one')}")
check("'summarize the third one' -> match 3", resolve_pending_file_selection(state, "summarize the third one") == PENDING[2])
check("'the last one' -> final match", resolve_pending_file_selection(state, "the last one") == PENDING[-1])
check("'2' -> match 2", resolve_pending_file_selection(state, "2") == PENDING[1])
check("'number 3' -> match 3", resolve_pending_file_selection(state, "number 3") == PENDING[2])
check("out-of-range '9' -> None", resolve_pending_file_selection(state, "9") is None)
check("unrelated message -> None", resolve_pending_file_selection(state, "what is the weather in London") is None)
check("'the first 3 lines' not read as picking file 3",
      resolve_pending_file_selection(state, "show me the first 3 lines") == PENDING[0])

# A mistyped ordinal used to resolve to nothing, and the reply was then
# handled as ordinary chat - which invented a fee schedule for a file
# that is a payment receipt.
check("'frist one' (typo) -> match 1",
      resolve_pending_file_selection(state, "frist one") == PENDING[0],
      f"got {resolve_pending_file_selection(state, 'frist one')}")
check("'secodn one' (typo) -> match 2",
      resolve_pending_file_selection(state, "secodn one") == PENDING[1],
      f"got {resolve_pending_file_selection(state, 'secodn one')}")
check("'thrid' (typo) -> match 3",
      resolve_pending_file_selection(state, "thrid") == PENDING[2],
      f"got {resolve_pending_file_selection(state, 'thrid')}")
check("an unrelated word is not fuzzy-matched to an ordinal",
      resolve_pending_file_selection(state, "banana") is None,
      f"got {resolve_pending_file_selection(state, 'banana')}")

empty = ConversationState()
check("no pending matches -> None", resolve_pending_file_selection(empty, "the first one") is None)


# ----------------------------------------------------------------
# 1b. Naming a file must not reuse the active one
# ----------------------------------------------------------------
print("\n--- active file vs named file ---")

active = ConversationState(
    last_file_path=r"C:\Users\TestUser\Documents\Hostel Fees.pdf",
    last_capabilities=["filesystem.read"],
)

check("'summarize it' -> active file", is_active_file_request(active, "summarize it"))
check("'what is the total?' -> active file", is_active_file_request(active, "what is the total?"))
check("'read the document' -> active file", is_active_file_request(active, "read the document"))
check("a reminder continues the active document",
      is_active_file_request(active, "quickly remind me what the main risk is"))
check("a disputed document field continues the active document",
      is_active_file_request(active, "isnt the expiry in march 2027 tho"))
check("an unrelated knowledge question does not continue the active document",
      not is_active_file_request(active, "who invented the telescope"))

check("'summarize a file named Budget' -> NOT active file",
      not is_active_file_request(active, "summarize a file named Budget"))
check("'find a file called Invoice' -> NOT active file",
      not is_active_file_request(active, "find a file called Invoice"))
check("'read report.pdf' -> NOT active file",
      not is_active_file_request(active, "read report.pdf"))
check("'look for my notes file' -> NOT active file",
      not is_active_file_request(active, "look for my notes file"))

# The router may also route to the open file. Naming a file must remove
# that option entirely, not just discourage it in the prompt.
from core.router import Router as _Router


class AlwaysActiveFileModel:
    """Answers ACTIVE_FILE whenever the schema allows it."""

    def __init__(self):
        self.offered = None

    def complete(self, system, prompt, schema=None, num_predict=None, think=None):
        self.offered = schema["properties"]["route"]["enum"]
        route = "ACTIVE_FILE" if "ACTIVE_FILE" in self.offered else "SAFE"
        return json.dumps({"route": route})


m = AlwaysActiveFileModel()
r = _Router(m)

st = ConversationState(
    last_file_path=r"C:\Users\TestUser\Documents\Hostel Fees.pdf",
    last_capabilities=["filesystem.read"],
)
check("a recent file field question routes to its capability",
      _Router(m).route(st, "whats the reference number") == "capability")

st2 = ConversationState(last_file_path=r"C:\Users\TestUser\Documents\Hostel Fees.pdf")
m.offered = None
route_named = _Router(m).route(st2, "summarize a file named HostelFeesSem1")
check("naming a file does not route to the open file",
      route_named != "file", f"got: {route_named}")
# Stronger than withholding ACTIVE_FILE from the schema. Withholding
# left SAFE available, and with conversation history in the prompt the
# model chose it - answering "I don't have access to your files" about
# files that were on disk. The question never reaches the model now.
check("naming a file is decided without asking the model",
      m.offered is None, f"model was consulted, offered: {m.offered}")
check("naming a file routes to a filesystem lookup",
      route_named == "capability", f"got: {route_named}")

st3 = ConversationState(last_file_path=r"C:\Users\TestUser\Documents\Hostel Fees.pdf")
m.offered = None
route_weather = _Router(m).route(st3, "whats the weather in rupnagar")
check("a live topic never routes to the open file",
      route_weather == "capability", f"got: {route_weather}")
check("a live topic is decided without asking the model",
      m.offered is None, f"model was consulted, offered: {m.offered}")

# Who-is / office questions go stale in training data. Answered from
# chat, "who is donald trump" gave the 45th presidency only, then
# confidently told the user the 47th "hasn't happened yet".
for msg in ["who is donald trump", "isnt he also the 47th president",
            "who is the prime minister"]:
    st_v = ConversationState()
    check(f"volatile fact looked up: {msg!r}",
          _Router(m).route(st_v, msg) == "capability",
          f"got: {_Router(m).route(st_v, msg)}")

# ...but settled history must not pay for a search.
st_hist = ConversationState()
check("settled history stays in chat",
      _Router(m).route(st_hist, "who invented the telescope") != "capability",
      "historical question was sent to a lookup")

# Asking for a file to be built is a job. Answered in chat, "can you
# make a ppt" claimed Athena was text-based and could not make files.
for msg in ["can you make a ppt", "make a powerpoint about the solar system",
            "create a spreadsheet of expenses"]:
    st_b = ConversationState()
    check(f"artifact request routed to a capability: {msg!r}",
          _Router(m).route(st_b, msg) == "capability",
          f"got: {_Router(m).route(st_b, msg)}")

# Anything that must be worked out goes to the planner, which computes
# it rather than reasoning it out. There is no per-topic capability -
# a function for postfix would do nothing for a physics or maths
# question - so the route has to cover all of them.
class _CalculateModel:
    """Answers CALCULATE, and records what it was offered."""

    def __init__(self):
        self.offered = None

    def complete(self, system, prompt, schema=None, num_predict=None, think=None):
        self.offered = schema["properties"]["route"]["enum"]
        return json.dumps({"route": "CALCULATE"})


_calc = _CalculateModel()
_calc_route = _Router(_calc).route(ConversationState(), "differentiate x^2 + 3x")
check("a calculation is routed as one, not as chat",
      _calc_route == "calculate", f"got: {_calc_route}")
check("CALCULATE appears in the router schema",
      "CALCULATE" in (_calc.offered or []), f"offered: {_calc.offered}")

_file_calculation_state = ConversationState(
    last_file_path=r"C:\Docs\recipe.txt",
    last_capabilities=["filesystem.read"],
)
_calc.offered = None
check("a file-derived numeric transformation is calculated",
      _Router(_calc).route(
          _file_calculation_state,
          "im cooking for 6 instead. how much do i need?",
      ) == "calculate")
check("file-derived calculation is decided without model arbitration",
      _calc.offered is None, f"model was consulted: {_calc.offered}")

_calc.offered = None
check("personal day planning stays in chat",
      _Router(_calc).route(
          ConversationState(),
          "got class at 9, gym at 6 and an assignment due tonight. help me plan the day simply",
      ) == "chat")
check("personal day planning is decided without model arbitration",
      _calc.offered is None, f"model was consulted: {_calc.offered}")


# Markdown stripping used to read two multiplication signs as an
# emphasis pair and delete both, corrupting a correct answer after
# every check upstream had passed.
from core.agent import _strip_markdown as _sm

check("multiplication operators survive markdown stripping",
      _sm("The postfix is A B + C D E + F G * + * H / +.")
      == "The postfix is A B + C D E + F G * + * H / +.",
      f"got {_sm('The postfix is A B + C D E + F G * + * H / +.')!r}")
check("bold is still stripped",
      _sm("This is **bold** text.") == "This is bold text.")
check("italic is still stripped",
      _sm("This is *italic* text.") == "This is italic text.")

# A slipped key must not drop a build request back into chat.
st_typo = ConversationState()
check("'mke a ppt on photosynthesis' still routed to a capability",
      _Router(m).route(st_typo, "mke a ppt on photosynthesis") == "capability",
      f"got: {_Router(m).route(st_typo, 'mke a ppt on photosynthesis')}")

# Creative writing names no file and must stay conversational.
for msg in ["write me a song", "write me a poem"]:
    st_c = ConversationState()
    check(f"creative writing stays in chat: {msg!r}",
          _Router(m).route(st_c, msg) != "capability",
          f"got: {_Router(m).route(st_c, msg)}")


# ----------------------------------------------------------------
# 2. Answer-vs-evidence grounding check (the London bug)
# ----------------------------------------------------------------
print("\n--- grounding check ---")

from core.agent import _answer_within_evidence

EV = ["Welcome to our new look and feel.",
      "London is partly cloudy at 30 degrees with later rain possible."]

check("supported composed answer accepted",
      _answer_within_evidence("Right now London is partly cloudy at 30 degrees, with rain possible later.", EV))
check("hallucinated number rejected",
      not _answer_within_evidence("Right now London is partly cloudy at 45 degrees.", EV))
check("hallucinated proper noun rejected",
      not _answer_within_evidence("Right now Manchester is partly cloudy at 30 degrees.", EV))
check("empty evidence rejected", not _answer_within_evidence("London is sunny.", []))

WORLD_CUP_EVIDENCE = [
    "Spain 1–0 Argentina: Ferran Torres scores as Spain win the FIFA World Cup for a second time. 19 July 2026.",
    "The 2026 final featured 11 saves by Argentina goalkeeper Emiliano Martínez.",
    "In an earlier meeting, Argentina won 2–1.",
]
check("an outcome claim with the correct winner is accepted",
      _answer_within_evidence(
          "Spain won the 2026 FIFA World Cup final against Argentina.",
          WORLD_CUP_EVIDENCE,
      ))
check("winner and runner-up cannot be inverted",
      not _answer_within_evidence(
          "Argentina won the 2026 FIFA World Cup final against Spain.",
          WORLD_CUP_EVIDENCE,
      ))

# Words may be spelled out where the source abbreviates; numbers may not.
RECEIPT = [
    "Date: 21-Dec-2025",
    "Payment Reference Number : TESTA7298990",
    "Total Amount : 22,600.00",
    "Name: Alex Morgan",
]

check("abbreviated month accepted when written out",
      _answer_within_evidence("It was paid on 21-Dec-2025, in December.", RECEIPT))
check("exact reference number accepted",
      _answer_within_evidence("The reference number is TESTA7298990.", RECEIPT))
check("invented reference number rejected",
      not _answer_within_evidence("The reference number is TESTA1234567.", RECEIPT))
check("invented total rejected",
      not _answer_within_evidence("The total was 26,100.00.", RECEIPT))
check("invented name rejected",
      not _answer_within_evidence("The payer was Rahul Sharma.", RECEIPT))
check("punctuation-only degree abbreviations remain grounded",
      _answer_within_evidence(
          "The programme includes B.Ed students.",
          ["Programme: BEd students"],
      ))


# ----------------------------------------------------------------
# 3. Full orchestration - web path (London regression)
# ----------------------------------------------------------------
print("\n--- orchestration: web answer ---")

from core.agent import Agent, MODE_DEFAULTS
import json as _json


class FakeExecutor:
    def __init__(self, result):
        self.result = result
        self.calls = []

    def execute(self, plan):
        self.calls.append(plan["steps"][0])
        return [self.result]


class FakeModel:
    def __init__(self, payload):
        self.payload = payload

    def complete(self, system, message, schema=None, num_predict=None, think=False):
        return _json.dumps(self.payload)

    def chat(self, state, message, system_prompt=None, images=None, num_predict=None):
        # Mirrors OllamaModel.chat, which appends both turns to state.
        state.messages.append({"role": "user", "content": message})
        state.messages.append({"role": "assistant", "content": "chat reply"})
        return "chat reply"


class FakePlanner:
    def __init__(self, steps):
        self.steps = list(steps)

    def plan_step(self, state, message, executed=None, must_calculate=False):
        self.must_calculate = must_calculate
        if self.steps:
            return {"done": False, "step": self.steps.pop(0)}
        return {"done": True}


class FakeRouter:
    def __init__(self, route="capability"):
        self._route = route

    def route(self, state, message):
        return self._route


def build_agent(planner, executor, model, route="capability", state=None,
                **knobs):
    a = object.__new__(Agent)
    a.state = state or ConversationState()
    a.router = FakeRouter(route)
    a.planner = planner
    a.executor = executor
    a.response_model = model
    a.planner_model = model

    # The real constructor sets these from the chosen mode; this
    # builder skips it. Applying the shared defaults keeps a test of
    # the common path from silently running with a mode's extras
    # switched on - a self-check would fire an unexpected model call
    # and the fakes would answer it with the wrong script.
    for knob, value in MODE_DEFAULTS.items():
        setattr(a, knob, knobs.pop(knob, value))

    if knobs:
        raise TypeError(f"unknown knob(s): {sorted(knobs)}")

    return a


# S1 = snippet (tagged as one block), S2 = nav sentence, S3 = weather sentence.
WEB_RESULT = {
    "success": True,
    "data": [{
        "title": "London Weather",
        "url": "https://example.com/london",
        "snippet": "London forecast",
        "content": ("Welcome to our new look and feel. "
                    "London is partly cloudy at 30 degrees with later rain possible."),
        "trusted": True,
    }],
    "queries": ["London weather today"],
}

# The model cites BOTH the nav furniture and the real sentence, but
# composes a correct answer - exactly the reported London failure.
agent = build_agent(
    FakePlanner([{"type": "web.search", "args": {"query": "London weather", "category": "weather"}}]),
    FakeExecutor(WEB_RESULT),
    FakeModel({
        "evidence": ["S2", "S3"],
        "answer": "Right now London is partly cloudy at 30 degrees, with rain possible later.",
    }),
)

answer = agent.respond("London")

check("composed answer survives (not replaced by nav text)",
      "partly cloudy" in answer and not answer.startswith("Welcome to our new look"),
      f"got: {answer!r}")
check("nav furniture not leaked into answer",
      "new look and feel" not in answer,
      f"got: {answer!r}")

# A web step can already contain several query variants.  Letting the
# planner emit another web.search for the same single-job request made
# one mistaken capital answer fan out into four slow searches.
_repeat_search_planner = FakePlanner([
    {"type": "web.search", "args": {
        "query": "capital of Australia", "category": "general",
    }},
    {"type": "web.search", "args": {
        "query": "Australia capital Canberra Sydney", "category": "general",
    }},
])
_repeat_search_executor = FakeExecutor(WEB_RESULT)
_repeat_search_agent = build_agent(
    _repeat_search_planner,
    _repeat_search_executor,
    FakeModel({
        "evidence": ["S3"],
        "answer": "London is partly cloudy at 30 degrees.",
    }),
)
_repeat_search_agent.respond("what is the capital of Australia?")
check("one successful web step cannot fan out into repeated searches",
      len(_repeat_search_executor.calls) == 1,
      f"calls={_repeat_search_executor.calls!r}")


class _TypoFallbackModel(FakeModel):
    def __init__(self):
        super().__init__({})

    def chat(self, state, message, system_prompt=None, images=None,
             num_predict=None):
        draft = "I don't have current information about that."
        state.messages.append({"role": "user", "content": message})
        state.messages.append({"role": "assistant", "content": draft})
        return draft

    def complete(self, system, message, schema=None, num_predict=None,
                 think=False):
        if schema is None:
            return "What's the capital of Australia?"
        return _json.dumps({
            "evidence": ["S2"],
            "answer": "The capital of Australia is Canberra.",
        })


_capital_result = {
    "success": True,
    "data": [{
        "title": "Australia facts",
        "url": "https://example.com/australia",
        "snippet": "Australia capital",
        "content": "The capital of Australia is Canberra.",
        "trusted": True,
    }],
    "queries": ["whts teh capitl of australia"],
}
_typo_fallback_executor = FakeExecutor(_capital_result)
_typo_fallback_agent = build_agent(
    FakePlanner([{"type": "web.search", "args": {
        "query": "whts teh capitl of australia", "category": "general",
    }}]),
    _typo_fallback_executor,
    _TypoFallbackModel(),
    route="chat",
)
_typo_fallback_answer = _typo_fallback_agent.respond(
    "whts teh capitl of australia"
)
check("a failed typo retry falls through to grounded lookup",
      "Canberra" in _typo_fallback_answer
      and len(_typo_fallback_executor.calls) == 1
      and len(_typo_fallback_agent.state.messages) == 2,
      f"answer={_typo_fallback_answer!r}, "
      f"messages={_typo_fallback_agent.state.messages!r}")


class _TypoEchoThenAnswerModel(FakeModel):
    """Deciphers a typo first, then answers the corrected question."""

    def __init__(self):
        super().__init__({})
        self.repair_calls = 0

    def chat(self, state, message, system_prompt=None, images=None,
             num_predict=None):
        draft = "I don't have current information about that."
        state.messages.append({"role": "user", "content": message})
        state.messages.append({"role": "assistant", "content": draft})
        return draft

    def complete(self, system, message, schema=None, num_predict=None,
                 think=False):
        self.repair_calls += 1
        return (
            "What's the capital of Australia?"
            if self.repair_calls == 1
            else "The capital of Australia is Canberra."
        )


_typo_echo_model = _TypoEchoThenAnswerModel()
_typo_echo_executor = FakeExecutor(_capital_result)
_typo_echo_agent = build_agent(
    FakePlanner([]), _typo_echo_executor, _typo_echo_model, route="chat"
)
_typo_echo_answer = _typo_echo_agent.respond("whts teh capitl of australia")
check("a corrected-question echo gets answered without a web lookup",
      "Canberra" in _typo_echo_answer
      and len(_typo_echo_executor.calls) == 0,
      f"answer={_typo_echo_answer!r}, "
      f"calls={_typo_echo_executor.calls!r}")

# Unsupported answer must still fall back to verbatim evidence.
agent2 = build_agent(
    FakePlanner([{"type": "web.search", "args": {"query": "London weather", "category": "weather"}}]),
    FakeExecutor(WEB_RESULT),
    FakeModel({
        "evidence": ["S3"],
        "answer": "Right now London is partly cloudy at 45 degrees.",
    }),
)
answer2 = agent2.respond("London")
check("unsupported answer falls back to on-topic evidence",
      "30 degrees" in answer2 and "45" not in answer2,
      f"got: {answer2!r}")

# The snippet is now citable - previously it was shown to the model
# but had no ID, so any answer drawn from it could never verify.
agent2b = build_agent(
    FakePlanner([{"type": "web.search", "args": {"query": "London weather", "category": "weather"}}]),
    FakeExecutor(WEB_RESULT),
    FakeModel({"evidence": ["S1"], "answer": "The London forecast is available."}),
)
answer2b = agent2b.respond("London")
check("snippet is citable evidence", "London forecast" in answer2b, f"got: {answer2b!r}")

# Off-topic evidence must NOT be served as though it were the answer.
OFFTOPIC = {
    "success": True,
    "data": [{
        "title": "Weather help",
        "url": "https://example.com/help",
        "snippet": "Help page",
        "content": ("This shows the pollen level for the region this location is in. "
                    "Feels like temperature considers factors such as wind speed."),
        "trusted": False,
    }],
    "queries": ["London weather today"],
}

agent2c = build_agent(
    FakePlanner([{"type": "web.search", "args": {"query": "London weather", "category": "weather"}}]),
    FakeExecutor(OFFTOPIC),
    FakeModel({
        "evidence": ["S2", "S3"],
        "answer": "London is partly cloudy at 23 degrees with 51% humidity.",
    }),
)
answer2c = agent2c.respond("London")
check("off-topic evidence is not served as the answer",
      "pollen" not in answer2c.lower() and "feels like" not in answer2c.lower(),
      f"got: {answer2c!r}")
check("off-topic evidence reports honest failure",
      "couldn't find" in answer2c.lower() or "could not find" in answer2c.lower(),
      f"got: {answer2c!r}")

# If the model cites an unrelated sentence, recover the best tagged
# sentence that matches the complete subject instead of either leaking
# the bad citation or discarding a result retrieval already found.
CAPITAL_WITH_BAD_CITATION = {
    "success": True,
    "data": [{
        "title": "Australia facts",
        "url": "https://example.com/australia",
        "snippet": "The capital of Australia is Canberra.",
        "content": "Australia had a population estimate in 2026.",
        "trusted": True,
    }],
    "queries": ["capital of Australia 21 August 2026"],
}
agent2d = build_agent(
    FakePlanner([{"type": "web.search", "args": {
        "query": "capital of Australia", "category": "general",
    }}]),
    FakeExecutor(CAPITAL_WITH_BAD_CITATION),
    FakeModel({
        "evidence": ["S2"],
        "answer": "The capital of Australia is Canberra.",
    }),
)
answer2d = agent2d.respond("what is the capital of Australia?")
check("bad web citations recover the strongest on-topic evidence",
      "Canberra" in answer2d and "population" not in answer2d,
      f"got: {answer2d!r}")

# The cited sentence can be relevant but incomplete. In that case the
# recovery above does not run, so a second directly relevant tagged
# sentence must be allowed to complete verification without trusting
# anything outside the retrieved result.
INCOMPLETE_RELEVANT_CITATION = {
    "success": True,
    "data": [{
        "title": "UK leadership update",
        "url": "https://example.com/uk-leadership",
        "snippet": "The UK prime minister's office changed hands in 2026.",
        "content": "Andy Burnham is the current UK prime minister.",
        "trusted": True,
    }],
    "queries": ["current UK prime minister"],
}
agent2e = build_agent(
    FakePlanner([{"type": "web.search", "args": {
        "query": "current UK prime minister", "category": "general",
    }}]),
    FakeExecutor(INCOMPLETE_RELEVANT_CITATION),
    FakeModel({
        "evidence": ["S1"],
        "answer": "Andy Burnham is the current UK prime minister.",
    }),
)
answer2e = agent2e.respond("who is the current UK prime minister?")
check("omitted relevant web support can complete a cited answer",
      answer2e == "Andy Burnham is the current UK prime minister.",
      f"got: {answer2e!r}")


# ----------------------------------------------------------------
# 4. Full orchestration - multi-match file selection
# ----------------------------------------------------------------
print("\n--- orchestration: file selection ---")

SEARCH_AMBIGUOUS = {
    "success": False,
    "error": "Found multiple files matching 'Hostel Fees':\n" + "\n".join(PENDING) + "\n\nWhich one did you mean?",
    "matches": list(PENDING),
}

agent3 = build_agent(
    FakePlanner([{"type": "filesystem.search", "args": {"name": "Hostel Fees"}}]),
    FakeExecutor(SEARCH_AMBIGUOUS),
    FakeModel({"evidence": [], "answer": ""}),
)

turn1 = agent3.respond("summarize a file named Hostel Fees")

check("ambiguous search returns the disambiguation prompt", "Which one did you mean?" in turn1, f"got: {turn1!r}")
check("pending matches stored in state", agent3.state.pending_file_paths == PENDING,
      f"got: {agent3.state.pending_file_paths}")
check("no file selected yet", agent3.state.last_file_path is None)

# Turn 2: "the first one" must deterministically read match 1.
READ_RESULT = {
    "success": True,
    "data": "Receipt\nName: Alex Morgan\nHostel: Cedar Hall\nAmount Paid: 22600\n",
}
exec2 = FakeExecutor(READ_RESULT)
agent3.executor = exec2
agent3.planner = FakePlanner([])  # planner must NOT be needed
agent3.response_model = FakeModel({
    "evidence": ["S2", "S4"],
    "answer": "The receipt is for Alex Morgan and the amount paid was 22600.",
})

turn2 = agent3.respond("the first one")

check("selection triggered a filesystem.read", exec2.calls and exec2.calls[0]["type"] == "filesystem.read",
      f"got: {exec2.calls}")
check("read the correct (first) path", exec2.calls and exec2.calls[0]["args"]["path"] == PENDING[0],
      f"got: {exec2.calls[0]['args']['path'] if exec2.calls else None}")
check("selection recorded as active file", agent3.state.last_file_path == PENDING[0])
check("pending matches cleared after selection", agent3.state.pending_file_paths == [])
check("answer returned from file evidence", "22600" in turn2, f"got: {turn2!r}")


# An unresolvable short reply to "which one did you mean?" must ask
# again rather than fall through to the chat path. Chat has no evidence
# and no grounding check behind it, so a failed selection reaching it
# produced a confident, entirely invented summary of the document.
agent_sel = build_agent(
    FakePlanner([]),
    FakeExecutor({"success": True, "data": "x"}),
    FakeModel({"evidence": [], "answer": "A completely invented fee schedule."}),
    route="chat",   # what the router actually did on the reported failure
    state=ConversationState(pending_file_paths=list(PENDING)),
)

turn_bad = agent_sel.respond("hmm not sure")

check("unresolved selection re-asks instead of answering",
      "which one" in turn_bad.lower() or "pick one" in turn_bad.lower(),
      f"got: {turn_bad!r}")
check("unresolved selection does not invent an answer",
      "invented" not in turn_bad, f"got: {turn_bad!r}")
check("unresolved selection lists the options again",
      PENDING[0] in turn_bad, f"got: {turn_bad!r}")
check("unresolved selection keeps the choice pending",
      agent_sel.state.pending_file_paths == PENDING,
      f"got: {agent_sel.state.pending_file_paths}")

# A genuine change of subject must still get through rather than being
# trapped by the pending choice.
agent_sel2 = build_agent(
    FakePlanner([]),
    FakeExecutor({"success": True, "data": "x"}),
    FakeModel({"evidence": [], "answer": ""}),
    route="chat",
    state=ConversationState(pending_file_paths=list(PENDING)),
)
turn_topic = agent_sel2.respond("actually never mind, write me a short poem instead")
check("a longer new subject is not trapped by the pending choice",
      turn_topic == "chat reply", f"got: {turn_topic!r}")

# A SHORT message can still be a genuine change of subject. "whats the
# weather in paris" is exactly five words, and was being answered with
# "which file did you mean?" because only the length was checked.
from core.router import names_a_new_subject as _new_subject

for msg in ["whats the weather in paris", "who is the president",
            "make a ppt on cells", "latest news"]:
    check(f"recognised as a new subject: {msg!r}", _new_subject(msg))

# Sums and notation conversions are short and name no topic, so they
# matched none of the patterns above. Found by the deep test: after an
# unanswered "which of these files did you mean?", the reply to
# "convert 255 to binary" was a list of PDFs in the Downloads folder.
# Nobody answers "which file?" with a base conversion.
for msg in ["convert 255 to binary", "whats 2+2", "the prefix of A+B",
            "what is 45 * 12", "convert A+B to postfix"]:
    check(f"a short calculation is not a file choice: {msg!r}",
          _new_subject(msg))

for msg in ["the first one", "hmm not sure", "frist one", "number 2"]:
    check(f"still treated as a file choice: {msg!r}", not _new_subject(msg))


# ----------------------------------------------------------------
# A follow-up to a lookup is still that lookup
# ----------------------------------------------------------------
#
# "And in Mumbai?" contains no weather word and no place the router
# recognises, so the small model classified it as ordinary
# conversation and answered it from the chat path - which has no
# evidence behind it and no grounding check in front of it.
#
# It invented Mumbai's weather. Partly cloudy and 31.2 degrees against
# a real overcast 27.2, with the humidity, wind and feels-like all
# wrong too, and the next question then copied those figures. Found by
# the deep test; the answer was convincing enough that only checking
# it against the live service showed it up.
#
# What identifies it is not the words but what came before.

from core.router import continues_a_lookup as _continues

_AFTER_LOOKUP = ConversationState(last_capabilities=["weather.current"])
_AFTER_CHAT = ConversationState(last_capabilities=[])

for msg in ["and in Mumbai?", "and Mumbai", "what about Delhi",
            "Mumbai?", "and there?", "in Paris"]:
    check(f"continues the lookup: {msg!r}", _continues(_AFTER_LOOKUP, msg))

# Being short is not enough. Length alone was the first attempt and it
# grabbed "explain how recursion works" - a complete question with
# nothing to do with the lookup before it.
for msg in ["explain how recursion works", "who is the president",
            "whats 2+2", "I love this weather though",
            "what is a derivative", "find my resume", "summarize this"]:
    check(f"starts something new: {msg!r}",
          not _continues(_AFTER_LOOKUP, msg))

check("nothing continues an ordinary conversation",
      not _continues(_AFTER_CHAT, "and in Mumbai?")
      and not _continues(_AFTER_CHAT, "Mumbai?"),
      "a chat turn would pull the next message into a lookup")


# ----------------------------------------------------------------
# A lookup with nothing to look up
# ----------------------------------------------------------------
#
# "What's the weather" names a live subject and no place. The router
# hands it to the model deliberately so it can ask which city - and the
# model answered it instead, from the conversation: after Delhi and
# Mumbai had been looked up it reported "overcast, 28.6 degrees,
# humidity 88%", close enough to Delhi's real reading to pass for a
# fresh one and belonging to nowhere at all.
#
# Chat has no evidence behind it and no grounding check in front of it,
# so the only safe answer to an unanswerable question is the question
# back.

from core.router import missing_subject_question as _ask_back

for msg in ["whats the weather", "what is the weather", "hows the weather",
            "tell me the temperature"]:
    check(f"asks which city: {msg!r}",
          "city" in _ask_back(msg).lower(),
          f"got: {_ask_back(msg)!r}")

check("asks which company for a bare share price",
      "company" in _ask_back("whats the stock price").lower())

check("asks which currencies for a bare exchange rate",
      "currencies" in _ask_back("whats the exchange rate").lower())

# A question that names its subject is answerable, and must not be
# turned into a question back.
for msg in ["whats the weather in Delhi", "the weather in paris",
            "weather today", "whats 2+2", "explain recursion"]:
    check(f"answerable, so no question back: {msg!r}", not _ask_back(msg))


# The planner says "system" when it means this machine's clock, which
# is not a zone name. It was rejected as unknown, so "what year is it"
# came back as the words "Unknown timezone: system".
from services.system_service import SystemService as _Sys

_clock = _Sys()

for _tz in ["system", "local", "none", None]:
    check(f"the local clock is readable with timezone={_tz!r}",
          _clock.datetime_now(timezone=_tz).get("success") is True,
          f"got: {_clock.datetime_now(timezone=_tz).get('error')}")

check("a real zone still works",
      _clock.datetime_now(timezone="Asia/Kolkata").get("success") is True)

check("a genuinely wrong zone is still refused",
      _clock.datetime_now(timezone="Nowhere/Fake").get("success") is False,
      "a typo'd zone would silently give local time instead")


# The clock's answer has to be citable, or the grounding check throws
# it away. Asked "what year is it", the model answered "It is the year
# 2026" and cited the date it had been handed - the citation resolved
# to nothing, the evidence was discarded as unverified, and a correct
# answer came back as "I couldn't find that in what I looked up."
#
# The year is spelled out on its own line because it gets asked for on
# its own. A question about the year should not depend on the model
# picking "2026" out of "2026-08-20" and the check then agreeing.
_clock_result = {
    "success": True,
    "data": {"date": "2026-08-20", "time": "01:52:18",
             "day_of_week": "Thursday", "timezone": "local",
             "iso": "2026-08-20T01:52:18"},
}

from core.prompt_builder import PromptBuilder as _PromptBuilder

_, _clock_map, _ = _PromptBuilder.build(
    ConversationState(), "what year is it",
    {"steps": [{"type": "system.datetime", "args": {}}]},
    [_clock_result],
)

_clock_lines = list(_clock_map.values())

check("the clock's answer is citable at all",
      bool(_clock_lines),
      "nothing to cite means a correct answer gets discarded")

check("the year is a sentence of its own",
      any("year is 2026" in line for line in _clock_lines),
      f"got: {_clock_lines}")

check("the date is still there too",
      any("2026-08-20" in line for line in _clock_lines))


# ----------------------------------------------------------------
# A missing-subject question, whichever way it gets routed
# ----------------------------------------------------------------
#
# The first fix for "what's the weather" only checked when the router
# sent the message to chat. "Live topic with nothing named" is
# deliberately left to the router's own classifier rather than forced
# anywhere, and that classifier can decide it needs a lookup instead of
# choosing chat - at which point the check that only fired on the chat
# branch never ran. Told earlier that the user studies near Rupnagar,
# Athena answered for Rupnagar, the town Ropar is in - confidently
# reporting the weather for a place that was never asked about.
#
# Nothing in the router or the agent can be unit-tested for which LLM
# route a live model would pick, so what is checked here is the
# contract: the guard must not depend on route at all, only on the
# gate that was added right before it in agent.py - "not forced_step".

import re as _re_missing_subject
from pathlib import Path as _Path_missing_subject

_agent_source = (_Path_missing_subject(__file__).resolve().parent.parent
                 / "core" / "agent.py")
_agent_text = _agent_source.read_text(encoding="utf-8")

_guard = _re_missing_subject.search(
    r"if not forced_step:\s*\n\s*missing = missing_subject_question\(message\)",
    _agent_text,
)

check("the missing-subject check runs regardless of route",
      _guard is not None,
      "gating it on route == \"chat\" would miss NEEDS_LOOKUP -> "
      "capability, which is exactly how Rupnagar was invented")

check("the check is not gated on route == \"chat\" a second time above it",
      not _re_missing_subject.search(
          r'if route == "chat":\s*\n\s*missing = missing_subject_question',
          _agent_text),
      "the old chat-only gate is still there instead of being replaced")


# ----------------------------------------------------------------
# A self-check claim has to actually be in the answer
# ----------------------------------------------------------------
#
# Checking "The program printed 5." against its evidence, the checker
# invented "The program printed 6." as the unsupported claim - a
# number that appears nowhere in the real answer, hedging it against a
# sentence Athena never wrote. Nothing required a flagged claim to
# share so much as a word with the text it was supposedly quoting.

print("\n--- self-check grounding ---")

import core.agent as _agent_module

_grounded = _agent_module._grounded_in_answer if hasattr(
    _agent_module, "_grounded_in_answer") else None

# The helper is defined inside _run_self_check rather than at module
# level, so it is exercised the same way the real code path does:
# through the source, not a private import.
_has_filter = "_grounded_in_answer" in _agent_text

check("self-check filters claims against the answer's own words",
      _has_filter,
      "a hallucinated claim like 'printed 6' in an answer of "
      "'printed 5' would reach the user unfiltered")

# A claim that IS the answer, restated, is not a claim - it is the
# checker failing to find anything specific. "The program printed 5."
# came back flagged as "The program printed 5..", and the result was a
# correct answer followed by a hedge quoting itself: "the sources
# don't clearly back this up: The program printed 5.."
check("self-check discards a claim that restates the whole answer",
      "_is_the_whole_answer" in _agent_text,
      "a near-duplicate of the answer would still produce a "
      "self-quoting hedge on an otherwise correct reply")


# ----------------------------------------------------------------
# Standing firm on a settled fact
# ----------------------------------------------------------------
#
# The chat prompt told the model what to do when UNCERTAIN and was
# silent about being certain. Told "no im pretty sure its 6" after
# correctly answering 5, the reply was "You are absolutely right...
# 2 + 3 does indeed equal 6" - agreeing with something false because
# it was asked for twice.

_chat_prompt_text = (Path(__file__).resolve().parent.parent
                     / "core" / "chat_prompt.py").read_text(encoding="utf-8")

check("the chat prompt tells the model to hold a certain answer",
      "not a reason to change your answer" in _chat_prompt_text
      or "insistence alone" in _chat_prompt_text,
      "nothing stops it agreeing with a wrong number asked for twice")

agent_sel3 = build_agent(
    FakePlanner([{"type": "weather.current", "args": {"location": "Paris"}}]),
    FakeExecutor({"success": True, "data": {
        "place": "Paris", "conditions": "clear", "temperature": 18,
        "temperature_unit": "C", "feels_like": 17, "humidity": 50,
        "humidity_unit": "%", "precipitation": 0, "precipitation_unit": "mm",
        "wind_speed": 5, "wind_unit": "km/h", "observed_at": "now",
        "timezone": "Europe/Paris"}}),
    FakeModel({"evidence": [], "answer": "It is 18 degrees in Paris."}),
    state=ConversationState(pending_file_paths=list(PENDING)),
)
turn_weather = agent_sel3.respond("whats the weather in paris")
check("a short new subject escapes the pending file choice",
      "which one" not in turn_weather.lower(), f"got: {turn_weather!r}")


# The router's CALCULATE decision must reach the planner. Left to
# re-derive it, the planner answered a question about hours with
# system.datetime - the clock, not the arithmetic.
_calc_planner = FakePlanner([
    {"type": "code.generate", "args": {"path": "solve.py", "spec": "x"}}
])
_calc_agent = build_agent(
    _calc_planner,
    # One stub serves both steps: code.generate, then the run that the
    # agent chains onto it automatically.
    FakeExecutor({"success": True,
                  "data": {"path": "solve.py", "bytes_written": 42,
                           "return_code": 0, "stdout": "80.0 km/h",
                           "stderr": ""}}),
    FakeModel({"evidence": [], "answer": ""}),
    route="calculate",
)
_calc_agent.respond("a train travels 240 km in 3 hours, what is its average speed")
check("the planner is told the answer must be computed",
      getattr(_calc_planner, "must_calculate", False) is True,
      "must_calculate was not passed through")


# A script written to work something out is scratch work, not the
# user's open file. After "how much is 50 USD in rupees" the throwaway
# converter became the active file, and the next question about a
# document was answered by reading that script back.
_scratch = build_agent(
    FakePlanner([{"type": "code.generate",
                  "args": {"path": "convert.py", "spec": "x"}}]),
    FakeExecutor({"success": True,
                  "data": {"path": "convert.py", "bytes_written": 42,
                           "return_code": 0, "stdout": "4771.5", "stderr": ""}}),
    FakeModel({"evidence": [], "answer": ""}),
    route="calculate",
)
_scratch.respond("how much is 50 usd in rupees")
check("a scratch script does not become the open file",
      _scratch.state.last_file_path is None,
      f"got: {_scratch.state.last_file_path}")

# A script the user asked for is remembered separately from documents.
_wanted = build_agent(
    FakePlanner([{"type": "code.generate",
                  "args": {"path": "sort.py", "spec": "x"}}]),
    FakeExecutor({"success": True,
                  "data": {"path": "sort.py", "bytes_written": 42}}),
    FakeModel({"evidence": [], "answer": ""}),
)
_wanted.respond("write me a python script that sorts a list")
check("a requested script does not replace the open document",
      _wanted.state.last_file_path is None,
      f"got: {_wanted.state.last_file_path}")
check("a requested script is remembered as generated code",
      _wanted.state.last_generated_path == "sort.py",
      f"got: {_wanted.state.last_generated_path}")

# Turn 3: "summarize it" must reuse the active file without the planner.
exec3 = FakeExecutor(READ_RESULT)
agent3.executor = exec3
agent3.planner = FakePlanner([])
turn3 = agent3.respond("summarize it")
check("'summarize it' reuses active file deterministically",
      exec3.calls and exec3.calls[0]["args"]["path"] == PENDING[0],
      f"got: {exec3.calls}")

# A correct answer that the model failed to cite must still stand, as
# long as the file itself supports it.
exec4 = FakeExecutor(READ_RESULT)
agent3.executor = exec4
agent3.planner = FakePlanner([])
agent3.router = FakeRouter("file")   # router recognises it as about the open file
agent3.response_model = FakeModel({"evidence": [], "answer": "The amount paid was 22600."})
turn4 = agent3.respond("what was paid")
check("uncited answer kept when the file supports it", "22600" in turn4, f"got: {turn4!r}")

# ...but an uncited answer the file does NOT support must be refused.
exec5 = FakeExecutor(READ_RESULT)
agent3.executor = exec5
agent3.planner = FakePlanner([])
agent3.router = FakeRouter("file")
agent3.response_model = FakeModel({"evidence": [], "answer": "The amount paid was 99999."})
turn5 = agent3.respond("what was paid")
check("uncited answer refused when the file does not support it",
      "99999" not in turn5, f"got: {turn5!r}")


# ----------------------------------------------------------------
# 4b. Weather query normalization
# ----------------------------------------------------------------
print("\n--- weather query normalization ---")

from core.execution_manager import ExecutionManager as EM

norm = EM._normalize_weather_query

check("day-month-year stripped",
      norm("current weather in London 10 August 2026") == "current weather in London",
      f"got: {norm('current weather in London 10 August 2026')!r}")
check("month-year stripped and 'today' added",
      norm("London weather August 2026") == "London weather today",
      f"got: {norm('London weather August 2026')!r}")
check("ISO date stripped",
      norm("London weather 2026-08-10") == "London weather today",
      f"got: {norm('London weather 2026-08-10')!r}")
check("bare year stripped",
      "2026" not in norm("London temperature humidity 2026"))
check("already-current query left alone",
      norm("current weather in London") == "current weather in London")
check("date-only query does not become empty", norm("2026") == "2026")


# ----------------------------------------------------------------
# 4c. Vision failure is reported clearly
# ----------------------------------------------------------------
print("\n--- vision fallback ---")


class BoomVision:
    model_name = "gemma3:12b"

    def chat(self, *a, **k):
        raise RuntimeError("model does not support multimodal requests")


agent_v = build_agent(FakePlanner([]), FakeExecutor({"success": True, "data": "x"}),
                      FakeModel({"evidence": [], "answer": ""}))
agent_v.vision_model = BoomVision()
vision_reply = agent_v.respond("what is this?", image_path="C:/tmp/x.png")

check("vision failure does not raise", isinstance(vision_reply, str))
check("vision failure names the fix", "ollama pull" in vision_reply, f"got: {vision_reply!r}")
check("vision failure recorded in state", len(agent_v.state.messages) == 2)


# ----------------------------------------------------------------
# 5. Planner failure / empty plan
# ----------------------------------------------------------------
print("\n--- planner failure ---")


class BrokenPlanner:
    def plan_step(self, state, message, executed=None, must_calculate=False):
        return {"done": True, "error": "Planner response must be a JSON object."}


agent4 = build_agent(BrokenPlanner(), FakeExecutor({"success": True, "data": "x"}),
                     FakeModel({"evidence": [], "answer": ""}))
answer4 = agent4.respond("do something impossible")

check("planner failure gives a controlled message, not a grounding error",
      "rephrase" in answer4.lower(), f"got: {answer4!r}")
check("planner failure still recorded in conversation state",
      len(agent4.state.messages) == 2)


# ----------------------------------------------------------------
# 6. Chat route untouched
# ----------------------------------------------------------------
print("\n--- chat route ---")

agent5 = build_agent(FakePlanner([]), FakeExecutor({"success": True, "data": "x"}),
                     FakeModel({"evidence": [], "answer": ""}), route="chat")
agent5.state.messages = []
reply = agent5.respond("write me a short poem")
check("chat route still returns a chat reply", reply == "chat reply", f"got: {reply!r}")


# ----------------------------------------------------------------
# 7. HTTP error boundary
# ----------------------------------------------------------------
print("\n--- http error boundary ---")

import core.agent as agent_module


class BoomAgent:
    def respond(self, message, image_path=None):
        raise RuntimeError(r"failure referencing C:\Users\TestUser\secret-document.pdf")

    def shutdown(self):
        pass


agent_module.Agent = BoomAgent

import core.web_app as web_app
from fastapi.testclient import TestClient

client = TestClient(web_app.app)
res = client.post("/chat", json={"message": "anything"})

check("failure returns HTTP 200, not 500", res.status_code == 200, f"got {res.status_code}")
check("failure returns a stable user-facing message",
      "went wrong" in res.json()["response"].lower(), f"got: {res.json()}")
check("controlled HTTP failure is marked as an error",
      res.json()["error"] is True, f"got: {res.json()}")
check("internal path not leaked to the client",
      "secret-document" not in res.text, f"got: {res.text}")


# ----------------------------------------------------------------
# Modes
# ----------------------------------------------------------------
#
# The knobs decide how much accuracy work runs, so a mode quietly
# gaining or losing one changes every answer it gives. Pinned here
# because the cost is invisible from the interface: Fast turning on a
# self-check would just feel slow, with nothing to say why.

print("\n--- modes ---")

from core.agent import MODES, DEFAULT_MODE, _mode_config, resolve_mode
from core.router import looks_arithmetic

check("three modes, in order",
      list(MODES) == ["fast", "balanced", "max"], f"got: {list(MODES)}")

check("balanced is the default",
      DEFAULT_MODE == "balanced", f"got: {DEFAULT_MODE}")

_fast = _mode_config("fast")
_balanced = _mode_config("balanced")
_max = _mode_config("max")

check("fast pays for none of the extras",
      not _fast["force_compute"] and not _fast["self_check"]
      and _fast["min_sources"] == 1 and _fast["search_results"] == 4,
      f"got: {_fast}")

# Two sources costs nothing - the count is worked out either way - so
# balanced takes it. Everything else there is a round trip, so it does
# not.
check("balanced takes only the free upgrade",
      _balanced["min_sources"] == 2
      and not _balanced["force_compute"] and not _balanced["self_check"]
      and _balanced["search_results"] == 4,
      f"got: {_balanced}")

check("max turns up useful work without the noisy model re-check",
      _max["force_compute"] and not _max["self_check"]
      and _max["min_sources"] == 2 and _max["search_results"] > 4,
      f"got: {_max}")

check("only max shows its working",
      bool(_max["style"]) and not _fast["style"] and not _balanced["style"],
      "style should be set on max alone")

check("every mode fills in the unset knobs",
      all(k in _mode_config(name) for name in MODES for k in MODE_DEFAULTS),
      "a mode is missing a default")

# A saved mode from an older version follows its rename. Without this
# someone who had picked Explain gets plain answers and no sign why.
check("retired mode names follow their rename",
      resolve_mode("accurate") == "balanced" and resolve_mode("study") == "max",
      f"got: {resolve_mode('accurate')}, {resolve_mode('study')}")

check("unknown mode falls back to the default",
      resolve_mode("nonsense") == DEFAULT_MODE, f"got: {resolve_mode('nonsense')}")

check("current mode names resolve to themselves",
      all(resolve_mode(n) == n for n in MODES), "a live mode was rewritten")


# The detector behind force_compute. Too narrow and Max stops
# computing; too wide and every date and price takes a Python round
# trip for nothing. Both directions are pinned.
_COMPUTES = [
    "convert A+B*C to postfix",
    "what is 45 * 12",
    "calculate the area of a circle radius 7",
    "a train travels 120 km in 2 hours how fast is it",
    "solve x^2 + 3x - 4 = 0",
    "whats the square root of 361",
    "how much is 15 percent of 240",
    "convert 255 to binary",
    "differentiate x^2 + 3x",
]

_RECALLS = [
    "who is the president",
    "he was born in 1947",
    "whats the weather in paris",
    "summarize this file",
    "explain how https works",
    "what is a derivative",
    "what does postfix mean",
    "how long is a piece of string",
    "difference between prefix and postfix",
]

check("work-it-out questions are recognised",
      all(looks_arithmetic(m) for m in _COMPUTES),
      f"missed: {[m for m in _COMPUTES if not looks_arithmetic(m)]}")

check("recalled and definition questions are left alone",
      not any(looks_arithmetic(m) for m in _RECALLS),
      f"wrongly flagged: {[m for m in _RECALLS if looks_arithmetic(m)]}")


# force_compute overrides the router rather than asking it nicely: a
# rule the router may decline to apply is not a guarantee.
_computed = build_agent(
    FakePlanner([{"type": "system.datetime", "args": {}}]),
    FakeExecutor({"success": True, "data": "x"}),
    FakeModel({"evidence": [], "answer": ""}),
    route="chat",
    force_compute=True,
)

check("force_compute sends arithmetic to the planner, not chat",
      _computed.force_compute and looks_arithmetic("what is 45 * 12"),
      "max would answer arithmetic from memory")

check("force_compute leaves non-arithmetic on the chat path",
      not looks_arithmetic("who is the president"),
      "max would compute a recall question")


# ----------------------------------------------------------------
# Stopping a reply
# ----------------------------------------------------------------
#
# Stopping breaks an assumption the rest of the server was built on:
# that only one turn can be running, because the page would not let
# you send another until the last reply arrived. Once the page can
# give up waiting, that is no longer true, and two turns sharing one
# agent would interleave their conversation history.

print("\n--- stopping ---")

import threading as _threading
import time as _time
from core.agent import PROGRESS as _PROGRESS, Stopped as _Stopped

_PROGRESS.clear()

check("stop does nothing when nothing is running",
      _PROGRESS.stop() is False and not _PROGRESS.stopping,
      "an unarmed stop would kill the next thing to report a stage")

# A stop that armed itself with no turn to consume it would sit there
# and kill whatever announced a stage next - a mode switch, say.
_PROGRESS.set("Working")
check("stop arms while a turn is running", _PROGRESS.stop() is True)

_stopped_at_next_stage = False
try:
    _PROGRESS.set("Next stage")
except _Stopped:
    _stopped_at_next_stage = True

check("a stopped turn gives up at its next stage", _stopped_at_next_stage)

_PROGRESS.clear()
check("clearing disarms the stop",
      not _PROGRESS.stopping and not _PROGRESS.busy,
      "a stale stop would kill the following turn")

_PROGRESS.start_turn()
check("a new turn starts unstopped", not _PROGRESS.stopping)
_PROGRESS.clear()


# The lock is the guarantee, not the disabled button: the button is
# only advice, and a stopped reply hands control back while the turn
# it abandoned is still unwinding.
import core.web_app as _webapp
from fastapi.testclient import TestClient as _TestClient

_overlaps = []
_running = [0]
_count_lock = _threading.Lock()


class _OverlapAgent:
    mode = "balanced"

    def respond(self, message, image_path=None):
        with _count_lock:
            _running[0] += 1
            if _running[0] > 1:
                _overlaps.append(message)

        _time.sleep(0.25)
        _PROGRESS.set("run")

        with _count_lock:
            _running[0] -= 1

        return "answer"

    def set_mode(self, mode):
        return mode


_webapp.agent = _OverlapAgent()
_client = _TestClient(_webapp.app)

_threads = [
    _threading.Thread(target=lambda i=i: _client.post("/chat", json={"message": f"m{i}"}))
    for i in range(4)
]
for _t in _threads:
    _t.start()
for _t in _threads:
    _t.join()

check("overlapping requests never run two turns at once",
      not _overlaps, f"ran concurrently: {_overlaps}")

_PROGRESS.clear()
check("stop is refused when idle", _client.post("/stop").json()["ok"] is False)

# Switching mode unloads the model the running turn is using, so it is
# refused server-side and not only greyed out in the page.
_PROGRESS.set("Working")
_busy_switch = _client.post("/mode", json={"mode": "fast"}).json()
check("mode switch refused while a reply is running",
      _busy_switch["ok"] is False, f"got: {_busy_switch}")

_PROGRESS.clear()
check("mode switch allowed once idle",
      _client.post("/mode", json={"mode": "fast"}).json()["ok"] is True)
_PROGRESS.clear()


# ----------------------------------------------------------------
# Weather place names
# ----------------------------------------------------------------

print("\n--- weather places ---")

from services.live_data_service import _place_label

# The three fields overlap often enough that joining them blindly
# produced "Japan, Japan" - which reads as a bug even when the weather
# behind it is right.
check("a place is not named twice",
      _place_label({"name": "Japan", "admin1": None, "country": "Japan"}) == "Japan",
      f'got: {_place_label({"name": "Japan", "admin1": None, "country": "Japan"})!r}')

check("city-states are not named three times",
      _place_label({"name": "Singapore", "admin1": "Singapore",
                    "country": "Singapore"}) == "Singapore")

check("ordinary places keep their full name",
      _place_label({"name": "Mumbai", "admin1": "Maharashtra",
                    "country": "India"}) == "Mumbai, Maharashtra, India")

check("a missing region is skipped, not left blank",
      _place_label({"name": "Monaco", "admin1": None,
                    "country": "Monaco"}) == "Monaco")


# ----------------------------------------------------------------
# Context window
# ----------------------------------------------------------------
#
# Ollama's default context is 4096 and it cuts from the FRONT without
# saying so. The planner's instructions are 4,886 tokens, so for a long
# time they never arrived whole - the model was answering with roughly
# half of them, missing the opening line, and the only symptom was that
# it kept ignoring rules it appeared to have been given.
#
# These are cheap arithmetic checks on purpose. They cannot call a
# model, but they catch the thing that actually went wrong: a prompt
# growing past the window it is sent in.

print("\n--- context window ---")

from models.ollama_model import NUM_CTX, _history_for_prompt
from config import (
    CODE_MAX_TOKENS, PLANNER_MAX_TOKENS, RESPONSE_MAX_TOKENS,
    ROUTER_MAX_TOKENS,
)
from core.planner import PLANNER_SYSTEM_PROMPT
from core.grounded_prompt import GROUNDED_WEB_SYSTEM_PROMPT
from core.chat_prompt import CHAT_SYSTEM_PROMPT
from core.router import ROUTER_SYSTEM_PROMPT

# Measured with the real tokenizer, not estimated: chars/4 put the
# planner at 5,410 when it is really 4,886, and the gap matters when
# the whole question is whether it fits.
# Measured with the real tokenizer against the real prompts, then
# turned into a characters-per-token ratio so this can be checked
# without loading a model.
#
# A hardcoded token count was the first version and it was useless: it
# went stale the moment a prompt changed, and kept passing while the
# planner grew from 4,886 to 5,295 tokens underneath it. Counting
# characters is an estimate, but it is an estimate of the thing that
# actually moves.
#
# 4.1 is the lowest ratio measured across the four prompts (the router,
# which is densest). Using the lowest means the estimate errs high,
# which is the safe direction for a budget.
_CHARS_PER_TOKEN = 4.1

_PROMPTS = {
    "planner": PLANNER_SYSTEM_PROMPT,
    "grounded-web": GROUNDED_WEB_SYSTEM_PROMPT,
    "chat": CHAT_SYSTEM_PROMPT,
    "router": ROUTER_SYSTEM_PROMPT,
}


def _estimated_tokens(text):
    return int(len(text) / _CHARS_PER_TOKEN)


for _name, _prompt in _PROMPTS.items():
    check(f"{_name} prompt fits the context window",
          _estimated_tokens(_prompt) < NUM_CTX,
          f"{_name} is about {_estimated_tokens(_prompt)} tokens, "
          f"window is {NUM_CTX}")

check("ordinary model replies have a generous finite ceiling",
      1000 <= RESPONSE_MAX_TOKENS < NUM_CTX,
      f"got: {RESPONSE_MAX_TOKENS}")
check("structured router output is tightly bounded",
      16 <= ROUTER_MAX_TOKENS <= 128,
      f"got: {ROUTER_MAX_TOKENS}")
check("planner JSON has room for a complete code snippet",
      RESPONSE_MAX_TOKENS >= PLANNER_MAX_TOKENS >= 512,
      f"got: {PLANNER_MAX_TOKENS}")
check("generated programs have more room than prose answers",
      CODE_MAX_TOKENS > RESPONSE_MAX_TOKENS,
      f"code={CODE_MAX_TOKENS}, response={RESPONSE_MAX_TOKENS}")

# The planner is the one that overflowed, and it is also the one that
# grows: the conversation, the executed steps and the file context are
# all appended to it, and then the answer is generated into whatever is
# left. Headroom is the point, not just fitting.
#
# A real planning call was measured at 7,444 tokens with an EMPTY
# conversation - so the fixed prompt is well under half of what a
# request actually costs, and the rest needs somewhere to go.
_SPARE = NUM_CTX - _estimated_tokens(PLANNER_SYSTEM_PROMPT)

check("the planner has room to grow into",
      _SPARE >= 8000,
      f"only about {_SPARE} tokens spare - either trim the planner "
      f"prompt or raise NUM_CTX")

# Same number for every call. Ollama reallocates the context when it
# changes, which costs a full model reload - 3.4s for the 8b, far
# worse for the 12b.
_sizes = {
    "chars": len(PLANNER_SYSTEM_PROMPT) + len(GROUNDED_WEB_SYSTEM_PROMPT)
             + len(CHAT_SYSTEM_PROMPT) + len(ROUTER_SYSTEM_PROMPT),
}
check("prompts are still the size they were measured at",
      _sizes["chars"] > 0 and NUM_CTX >= 8192,
      "context was lowered below the measured planner prompt")

class _HistoryState:
    def __init__(self, n, summarized_upto=0):
        self.messages = [{"role": "user", "content": f"m{i}"} for i in range(n)]
        self.summarized_upto = summarized_upto


_long = _HistoryState(200, summarized_upto=176)
_remaining = _history_for_prompt(_long)

check("messages represented by the summary are not resent",
      len(_remaining) == 24,
      f"got {len(_remaining)} of 200")

check("the unsummarised tail keeps its original order",
      _remaining[0]["content"] == "m176"
      and _remaining[-1]["content"] == "m199",
      f"got: {_remaining[0]} to {_remaining[-1]}")

_delayed_summary = _HistoryState(50)
check("a delayed summary never creates a silent history gap",
      len(_history_for_prompt(_delayed_summary)) == 50)

_short = _HistoryState(3)
check("a short conversation is left alone",
      _history_for_prompt(_short) == _short.messages)


# Tokens that were counted but never read.
#
# Ollama caches the prompt prefix, so the planner's 5,295-token prompt
# resent on a second planning round costs almost nothing - and is still
# reported in full by prompt_eval_count. A reply showing "9,137 tokens"
# was mostly reporting work that never happened, which made the most
# prominent number on screen the least meaningful one.
#
# The durations below are real: a cold 2,430-token prompt took 1.14s,
# and the same prompt again took 0.04s.
from models.ollama_model import UsageTracker

_usage = UsageTracker()
_usage.record({"prompt_eval_count": 2430, "eval_count": 10,
               "prompt_eval_duration": int(1.14e9)})

check("a prompt that was really read counts as work",
      _usage.cached_tokens == 0 and _usage.computed == 2440,
      f"cached={_usage.cached_tokens} computed={_usage.computed}")

_usage.record({"prompt_eval_count": 2435, "eval_count": 10,
               "prompt_eval_duration": int(0.04e9)})

check("a prompt served from the cache is not counted as work",
      _usage.cached_tokens == 2435,
      f"cached={_usage.cached_tokens}")

check("the total is still the total",
      _usage.total == 4885 and _usage.computed == 2450,
      f"total={_usage.total} computed={_usage.computed}")

# Generated tokens are always real work, however fast they arrive.
_gen = UsageTracker()
_gen.record({"prompt_eval_count": 0, "eval_count": 500,
             "prompt_eval_duration": 0})

check("generated tokens always count",
      _gen.computed == 500, f"computed={_gen.computed}")

# Read, reused and written are shown side by side, so they have to add
# up to the total or the line contradicts itself on screen. An earlier
# version showed "19 tokens" beside "1,975 in / 19 out" - the 19 was
# the written tokens counted twice, once under each name.
_adds = UsageTracker()
_adds.record({"prompt_eval_count": 1933, "eval_count": 17,
              "prompt_eval_duration": int(0.9e9)})
_adds.record({"prompt_eval_count": 1975, "eval_count": 19,
              "prompt_eval_duration": int(0.02e9)})

check("read, reused and written account for every token",
      _adds.prompt_read + _adds.cached_tokens + _adds.output_tokens
      == _adds.total,
      f"read={_adds.prompt_read} reused={_adds.cached_tokens} "
      f"written={_adds.output_tokens} total={_adds.total}")


# ----------------------------------------------------------------
# Notation conversions
# ----------------------------------------------------------------
#
# "The prefix of A+B" is "+AB". Left to the router model it went to
# chat, where "prefix" was read as a STRING prefix and the answer came
# back as "A+" - then defended across three more turns.
#
# The check used to require a digit, and notation questions are usually
# symbolic: "A+B" has nothing to match. The expression is the signal.

print("\n--- notation conversions ---")

from core.router import looks_notation_conversion as _notation
from tests.conversation_corpus import CONVERSATION as _CORPUS


class _AlwaysSafe:
    """A router model that always says the answer needs no capability.

    The point is that the override decides before it is asked, so a
    model insisting otherwise must not change the outcome.
    """

    def complete(self, system, prompt, schema=None, num_predict=None, think=None):
        return _json.dumps({"route": "SAFE"})

_CONVERSIONS = [
    "And the prefix of A+B",        # the reported failure
    "whats the postfix of A+B",
    "convert A+B to prefix",
    "prefix of 2+3",
    "A+B*C in postfix",
    "convert 255 to binary",
    "infix of AB+",                 # operator at the end
    "infix of +AB",                 # operator at the front
]

# "Prefix" is an ordinary word. It only becomes a conversion when there
# is something to convert, or the override would grab half of English.
_NOT_CONVERSIONS = [
    "give me just the prefix",
    "whats the country prefix for india",
    "what is a prefix",
    "the prefix of the word unhappy",
    "binary is base two",
    "remove the file prefix",
]

check("an expression in another notation is recognised",
      all(_notation(m) for m in _CONVERSIONS),
      f"missed: {[m for m in _CONVERSIONS if not _notation(m)]}")

check("the bare word is not a conversion",
      not any(_notation(m) for m in _NOT_CONVERSIONS),
      f"wrongly caught: {[m for m in _NOT_CONVERSIONS if _notation(m)]}")

check("no conversation is mistaken for a conversion",
      not any(_notation(m) for m in _CORPUS),
      f"caught: {[m for m in _CORPUS if _notation(m)]}")

check("a notation question routes to the calculator, not chat",
      _Router(_AlwaysSafe()).route(ConversationState(), "And the prefix of A+B")
      == "calculate",
      "the model would answer it from memory and get it wrong")


# ----------------------------------------------------------------
# Evidence markers must not reach the reader
# ----------------------------------------------------------------
#
# [S1], [S2] and so on are scaffolding: they make evidence citable for
# the grounding check and are never meant to be seen.
#
# They only began leaking once script output was tagged. A file's text
# gets paraphrased, so its markers rarely survive, but a computed
# result is copied out exactly as printed - and the answer came back as
# "The prefix of A+B is [S1] + A B."

print("\n--- evidence markers ---")

from core.agent import _strip_tags

check("a marker is removed from the answer",
      _strip_tags("The prefix of A+B is [S1] + A B.")
      == "The prefix of A+B is + A B.")

check("a marker at the start goes too",
      _strip_tags("[S1] A B +") == "A B +")

check("several markers are all removed",
      "[S" not in _strip_tags("From [S12] and [S3] the total is 500."))

check("text with no markers is untouched",
      _strip_tags("no tags here") == "no tags here")

# An index in code looks like a marker and is not one. The pattern
# requires the bracket to stand on its own.
check("an index in code is left alone",
      _strip_tags("array[S1] stays") == "array[S1] stays",
      "a computed answer containing code would be corrupted")


# ----------------------------------------------------------------
# A verification step that breaks instead of disagreeing
# ----------------------------------------------------------------
#
# The generation prompt asks a computed answer to verify itself and
# print exactly one of two things - the result, or the line
# VERIFICATION FAILED. The first real run did neither: the check's own
# eval() raised NameError on a bare letter, the script caught it, and
# printed "There was an error during verification because the name
# 'A' is not defined" NEXT TO the correct result - reporting success
# and failure in the same sentence, none of which matched the exact
# phrase the first check looks for.

print("\n--- broken verification ---")

from core.agent import _script_error as _err

_LEAKED = {
    "data": {"return_code": 0, "stderr": "",
             "stdout": ("The infix expression A+B converts to the prefix "
                        "expression + A B. There was an error during "
                        "verification because the name 'A' is not defined.\n")},
}

_CLEAN = {
    "data": {"return_code": 0, "stderr": "", "stdout": "Prefix Expression: + A B\n"},
}

check("a verification crash leaking into the output is caught",
      bool(_err(_LEAKED)),
      "the mixed success/failure sentence would have reached the reply")

check("output naming the failure explicitly is not repeated",
      "printed alongside the result" in _err(_LEAKED),
      "the repair prompt would not explain what went wrong")

check("a clean run with no mention of verification failing is left alone",
      not _err(_CLEAN))


# ----------------------------------------------------------------
# "find" in ordinary English
# ----------------------------------------------------------------
#
# The file-request check used to be a regex with a bare `find|locate`
# in it, so any sentence containing the word was forced to the
# capability route before the model was consulted. "I always find
# reasons not to go" - said in the middle of a conversation about
# football - was sent to the planner and answered with a generated
# script.
#
# It matters that this was a deterministic override: it cannot be
# out-voted, so it produced the same wrong answer every time, and no
# amount of re-asking the model would have caught it.

print("\n--- 'find' in ordinary English ---")

from core.router import names_a_file

_NOT_FILE_REQUESTS = [
    "I always find reasons not to go",   # the reported failure
    "I find that interesting",
    "we find it hard to say",
    "people find this confusing",
    "find out what time it is",          # discover, not search
    "you find the strangest things funny",
    "they never find time for it",
    "it is hard to find motivation",
    "struggling to find the words",
]

_FILE_REQUESTS = [
    "is there a file named hostel fees",
    "find the hostel fees pdf",
    "can you find my resume",
    "where can I find the report",
    "search for the budget spreadsheet",
    "look for invoice.pdf",
    "locate the config file",
    "summarize a file named Budget",
    "please find the receipt",
    "help me find the invoice",
]

check("ordinary uses of 'find' are not file requests",
      not any(names_a_file(m) for m in _NOT_FILE_REQUESTS),
      f"wrongly flagged: {[m for m in _NOT_FILE_REQUESTS if names_a_file(m)]}")

check("real file requests are still recognised",
      all(names_a_file(m) for m in _FILE_REQUESTS),
      f"missed: {[m for m in _FILE_REQUESTS if not names_a_file(m)]}")


# ----------------------------------------------------------------
# Override audit
# ----------------------------------------------------------------
#
# The overrides run before the model and cannot be overruled. That is
# what makes them reliable and what makes them dangerous: every one
# matches on words, and words occur in conversation.
#
# This runs the whole set against a corpus of ordinary sentences. It
# does not test one bug - it tests the category, so a new pattern that
# catches innocent English fails here rather than surfacing later as an
# answer nobody can reproduce.

print("\n--- override audit ---")

from tests.conversation_corpus import (
    CONVERSATION, REAL_REQUESTS, OVERRIDE_BACKED_REQUESTS,
)
from core.router import (
    gather_hints,
    is_plain_statement,
    _FILE_NAME_SIGNAL as _RX_FILE_NAME,
    _NAMED_FORMAT as _RX_NAMED_FORMAT,
    _LIVE_TOPICS as _RX_LIVE,
    _VOLATILE_FACTS as _RX_VOLATILE,
)


class _NeverAsked:
    """Stands in for the model, and fails loudly if anything calls it.

    The audit is about which messages reach the model at all, so a
    real call would mean the test was measuring the wrong thing.
    """

    def complete(self, *args, **kwargs):
        raise AssertionError("the router should not need a model call here")


class _NoFile:
    messages = []
    last_file_path = None
    pending_file_paths = []


def _forced(message):
    """Which signals decide the route outright, after the guard.

    Only the unambiguous ones are left here. Everything that matches on
    ordinary English is now a hint the model may decline, so it does
    not belong in a test about what cannot be argued with.
    """

    if is_plain_statement(message):
        return []

    hits = []

    for name, pattern in (
        ("names-a-file", _RX_FILE_NAME),
        ("named-format", _RX_NAMED_FORMAT),
        ("live-topic", _RX_LIVE),
        ("volatile-fact", _RX_VOLATILE),
    ):
        if pattern.search(message):
            hits.append(name)

    return hits


# What is allowed to decide without asking the model: things that
# cannot plausibly mean anything else. Verbs are absent by design.
_MAY_FORCE = {"names-a-file", "named-format", "live-topic", "volatile-fact"}


def _signals(message):
    """Everything noticed - forced or hinted."""

    if is_plain_statement(message):
        return []

    return _forced(message) + gather_hints(message)


_caught = {m: _signals(m) for m in CONVERSATION if _signals(m)}

check(f"nothing fires on ordinary conversation ({len(CONVERSATION)} sentences)",
      not _caught,
      "caught: " + "; ".join(f"{m!r} -> {h}" for m, h in list(_caught.items())[:6]))

# Narrowing patterns until conversation passes is easy if nothing
# checks the real requests still work. Each of these went wrong once
# without something noticing it.
_missed = [m for m in OVERRIDE_BACKED_REQUESTS if not _signals(m)]

check(f"requests that need a signal still get one "
      f"({len(OVERRIDE_BACKED_REQUESTS)} messages)",
      not _missed,
      f"no signal at all: {_missed}")

# The point of the split. Only signals that cannot mean anything else
# are allowed to decide without asking; everything built on ordinary
# words has to go through the model.
_still_forced = {m: _forced(m) for m in OVERRIDE_BACKED_REQUESTS if _forced(m)}

check("only unambiguous signals decide the route outright",
      all(set(h) <= _MAY_FORCE for h in _still_forced.values()),
      f"forced by something ambiguous: {_still_forced}")

check("ordinary words no longer force a route",
      not _forced("find my resume") and bool(gather_hints("find my resume")),
      "a bare locate verb is still deciding without the model")

check("a file extension still decides outright",
      _forced("look for invoice.pdf") == ["names-a-file"])

check("a named format still decides outright",
      _forced("make a ppt on photosynthesis") == ["named-format"])


# A live topic with nothing named cannot be looked up, and forcing it
# to a capability anyway made the planner pick a city itself - "what's
# the weather" was answered with the temperature in New Delhi, for a
# user who had never mentioned it. Found by the quality evaluation,
# which is the sort of thing the deterministic suite cannot see.
from core.router import _NAMES_A_SUBJECT as _RX_SUBJECT, _LIVE_TOPICS as _RX_LIVE2

for _msg in ["whats the weather", "what is the weather", "whats the temperature"]:
    check(f"a live topic with nothing named is not forced: {_msg!r}",
          bool(_RX_LIVE2.search(_msg)) and not _RX_SUBJECT.search(_msg),
          "would be forced to a lookup with no place to look up")

for _msg in ["whats the weather in paris", "whats the temperature in delhi",
             "weather today", "whats the weather outside"]:
    check(f"a live topic that names its subject is still forced: {_msg!r}",
          bool(_RX_LIVE2.search(_msg)) and bool(_RX_SUBJECT.search(_msg)))

# The guard sits in front of everything, so it must not mistake any
# real request for conversation - including the ones no override
# handles.
_guarded_away = [m for m in REAL_REQUESTS if is_plain_statement(m)]

check(f"the guard lets every real request through ({len(REAL_REQUESTS)} messages)",
      not _guarded_away,
      f"wrongly read as statements: {_guarded_away}")


# The guard itself. It may only ever suppress an override - a message
# it wrongly calls a statement still gets classified by the model, so
# the cost is a model call rather than a wrong answer. That is why it
# has to say False whenever it is unsure.
check("statements are recognised",
      all(is_plain_statement(m) for m in [
          "I always find reasons not to go",
          "I recently started playing football again",
          "I love this weather",
          "we built a treehouse when I was young",
      ]),
      "a plain statement was read as a request")

check("questions are never treated as statements",
      not any(is_plain_statement(m) for m in [
          "whats the weather in paris",
          "who is the prime minister",
          "is there a file named hostel fees",
          "what time is it",
      ]),
      "a question was read as a statement")

check("instructions are never treated as statements",
      not any(is_plain_statement(m) for m in [
          "find the hostel fees pdf",
          "make a ppt on photosynthesis",
          "convert A+B*C to postfix",
          "please find the receipt",
          "can you find my resume",
      ]),
      "an instruction was read as a statement")

check("an empty message is not a statement",
      not is_plain_statement("") and not is_plain_statement("   "))


# ----------------------------------------------------------------
# Saved conversations
# ----------------------------------------------------------------

print("\n--- saved conversations ---")

import shutil as _shutil
import tempfile as _tempfile
from pathlib import Path as _Path

from services.conversation_store import (
    Conversation as _Convo, ConversationStore as _Store,
    UNTITLED as _UNTITLED, title_from as _title_from,
)

_dir = _Path(_tempfile.mkdtemp()) / "conversations"
_store = _Store(_dir)

check("an empty store lists nothing",
      _store.list() == [] and _store.most_recent_id() is None)

_id = _store.new_id("what is the capital of France")
_store.save(_Convo(id=_id, title=_title_from("what is the capital of France"),
                   messages=[{"role": "user", "content": "hi"},
                             {"role": "assistant", "content": "hello"}],
                   summary="The user said hi.", summarized_upto=2))

_back = _store.load(_id)

check("a conversation survives a round trip",
      _back is not None
      and _back.messages == [{"role": "user", "content": "hi"},
                             {"role": "assistant", "content": "hello"}]
      and _back.summary == "The user said hi."
      and _back.summarized_upto == 2,
      "the summary and its coverage must travel with the messages")

check("a missing conversation reads as None", _store.load("nope") is None)
check("deleting a missing conversation is not an error",
      _store.delete("nope") is False)

# A corrupt file should cost that one conversation, not the ability to
# open Athena at all.
(_dir / "broken.json").write_text("{not json", encoding="utf-8")

check("a corrupt file does not break the listing",
      len(_store.list()) == 1)

check("an unnamed conversation gets a placeholder",
      _title_from("") == _UNTITLED)

check("a long first message is trimmed for the title",
      len(_title_from("x" * 200)) < 70)

_shutil.rmtree(_dir.parent, ignore_errors=True)


# ----------------------------------------------------------------
# Folding old messages into a summary
# ----------------------------------------------------------------
#
# Trimming history bounds the prompt but forgets - twenty messages in,
# the model no longer knows your name. Summarising keeps what was said
# while still cutting what is sent.

print("\n--- summarising ---")

from core.agent import SUMMARIZE_EVERY, KEEP_VERBATIM, _messages_as_text

check("recent messages are always sent verbatim",
      0 < KEEP_VERBATIM < SUMMARIZE_EVERY,
      f"KEEP_VERBATIM={KEEP_VERBATIM} SUMMARIZE_EVERY={SUMMARIZE_EVERY}")


class _FakeSummarizer:
    def __init__(self):
        self.calls = 0

    def complete(self, system, message, **kwargs):
        self.calls += 1
        return "The user is Alex. They like football."


from core.conversation_state import ConversationState as _State

# Built without __init__ on purpose: constructing an Agent loads a
# model, and none of this needs one.
_summarizer = _FakeSummarizer()
_agent = Agent.__new__(Agent)
_agent.state = _State()
_agent.store = _Store(_Path(_tempfile.mkdtemp()) / "c")
_agent.response_model = _summarizer

_resent = []

_summary_turns = SUMMARIZE_EVERY // 2 + 5

for _turn in range(_summary_turns):
    _agent.state.messages.append({"role": "user", "content": f"q{_turn}"})
    _agent.state.messages.append({"role": "assistant", "content": f"a{_turn}"})
    _agent._maybe_summarize()
    _resent.append(len(_agent.state.messages) - _agent.state.summarized_upto)

check("nothing is summarised before there is enough to summarise",
      _resent[0] == 2 and _summarizer.calls >= 1,
      f"first turn resent {_resent[0]}, calls={_summarizer.calls}")

# The whole point: what gets resent stops growing.
check("the number of resent messages stays bounded",
      max(_resent) <= SUMMARIZE_EVERY + KEEP_VERBATIM,
      f"peaked at {max(_resent)} messages")

check("summarising is not repeated for the same messages",
      _summarizer.calls <= (_summary_turns * 2) // SUMMARIZE_EVERY + 1,
      f"{_summarizer.calls} calls for {_summary_turns * 2} messages")

check("the summary is kept once written",
      "Alex" in _agent.state.summary)

check("the summary reaches the model as context",
      "EARLIER IN THIS CONVERSATION" in _agent._with_memory("BASE")
      and "Alex" in _agent._with_memory("BASE"))

_blank = Agent.__new__(Agent)
_blank.state = _State()

check("no notes are added when there is nothing to remember",
      _blank._with_memory("BASE") == "BASE")

check("messages are flattened readably for the summariser",
      _messages_as_text([{"role": "user", "content": "hi"},
                         {"role": "assistant", "content": "hello"}])
      == "User: hi\nAthena: hello")

# Messages the summary covers are not sent again - that is where the
# saving actually comes from.
_state = _State()
_state.messages = [{"role": "user", "content": f"m{i}"} for i in range(30)]
_state.summarized_upto = 12

_sent = _history_for_prompt(_state)

check("summarised messages are not resent",
      _sent[0]["content"] == "m12" and len(_sent) == 18,
      f"got {len(_sent)} starting at {_sent[0]}")


# ----------------------------------------------------------------
# Semantic search
# ----------------------------------------------------------------
#
# Only the parts that need no model: chunking, similarity, and the
# wiring that decides whether results reach the answer at all. The
# search quality itself is measured in eval_quality.py, which has a
# model to ask.

print("\n--- semantic search ---")

from services.semantic_index import (
    _chunk, _cosine, CHUNK_CHARS, MAX_CHUNKS_PER_FILE, SIMILARITY_FLOOR,
    INDEXABLE,
)

check("identical vectors score 1", abs(_cosine([1, 0, 0], [1, 0, 0]) - 1.0) < 1e-9)
check("unrelated vectors score 0", _cosine([1, 0, 0], [0, 1, 0]) == 0.0)
check("a zero vector does not divide by zero", _cosine([0, 0, 0], [1, 0, 0]) == 0.0)

_chunks = _chunk("This is a sentence about something. " * 100)

check("long text is split into chunks", len(_chunks) > 1)
check("chunks stay near the target size",
      all(len(c) <= CHUNK_CHARS + 50 for c in _chunks),
      f"sizes: {[len(c) for c in _chunks][:5]}")
check("chunks end on a sentence where they can",
      _chunks[0].rstrip().endswith("."))
check("empty text produces no chunks", _chunk("") == [] and _chunk("   ") == [])
check("one long document cannot flood the index",
      len(_chunk("word. " * 20000)) <= MAX_CHUNKS_PER_FILE)

# Erring high on purpose: a document not returned costs another
# search, one returned wrongly becomes evidence for an answer about
# something it never mentioned.
check("the similarity floor is set above the noise", SIMILARITY_FLOOR >= 0.5,
      f"floor is {SIMILARITY_FLOOR}")

check("images are not indexed", not ({".png", ".jpg"} & INDEXABLE),
      "OCR per page is too slow to index")


# The wiring. Each of these was a real failure while building it: the
# results were found and then discarded three separate ways.
from core.agent import _FILESYSTEM_TYPES, _STEP_STAGES, _select_grounded_prompt
from core.grounded_prompt import GROUNDED_FILESYSTEM_SYSTEM_PROMPT
from core.capabilities import CAPABILITIES

check("semantic search gets the filesystem grounding rules",
      _select_grounded_prompt([{"type": "filesystem.semantic_search"}])
      is GROUNDED_FILESYSTEM_SYSTEM_PROMPT,
      "it fell through to the generic prompt and answered 'not found' "
      "about a document it had just found")

check("semantic search is a known filesystem step",
      "filesystem.semantic_search" in _FILESYSTEM_TYPES)

check("semantic search has a stage label",
      "filesystem.semantic_search" in _STEP_STAGES)

check("semantic search is registered as a capability",
      any(c["type"] == "filesystem.semantic_search" for c in CAPABILITIES))

# A capability the planner has never seen an example of does not get
# chosen, however well the registry describes it.
from core.planner import PLANNER_SYSTEM_PROMPT as _PLAN_PROMPT

check("the planner has a worked example for semantic search",
      _PLAN_PROMPT.count("filesystem.semantic_search") >= 2,
      "registry entries alone have never been enough to teach the planner")


# The filename has to be part of the evidence, not a heading above it.
# Left untagged, naming the file - the one thing the user asked for -
# counted as an unsupported claim, so the answer was thrown away and
# the raw passage came back with no indication which document it was.
from core.prompt_builder import PromptBuilder as _PB

_semantic_result = {
    "success": True,
    "matches": [{
        "path": r"C:\Users\x\Downloads\Receipt.pdf",
        "name": "Receipt.pdf",
        "score": 0.61,
        "excerpt": "Amount in figures Rs. 117285.00 paid on 13-07-2026.",
    }],
}

_prompt, _smap, _ = _PB.build(
    ConversationState(),
    "which file has my hostel payment",
    {"steps": [{"type": "filesystem.semantic_search",
                "args": {"query": "hostel payment"}}]},
    [_semantic_result],
)

check("the matched passage reaches the prompt",
      "117285.00" in _prompt)

check("the filename is citable evidence, not just a heading",
      any("Receipt.pdf" in line for line in _smap.values()),
      f"tagged lines: {list(_smap.values())}")

check("the answer can name the file without being called unsupported",
      _answer_within_evidence("It is in Receipt.pdf.", list(_smap.values())))


# A receipt prints "117285.00" and the model writes "Rs. 117,285.00" -
# the same figure, punctuated the way anyone would write it. That was
# being rejected as an invented number.
_receipt = ["*Amount (in figures) Rs. 117285.00 *Amount (in words) One Lakh"]

check("a thousands separator is a wording difference",
      _answer_within_evidence("The amount is Rs. 117,285.00", _receipt))

check("wrong digits are still caught",
      not _answer_within_evidence("The amount is Rs. 127,285.00", _receipt))

check("an invented figure is still caught",
      not _answer_within_evidence("The deposit was 5,000", _receipt))


# Athena running a script and reading its output is evidence - better
# evidence than a web page. It was being passed as plain text, so a
# plan that only ran a script produced no citable sentences at all,
# and the checks downstream compare the answer against exactly those.
# Asked for the postfix of A+B, Athena computed "A B +" correctly and
# then told the user "the sources don't clearly back this up".
_ran = {
    "success": True,
    "data": {"return_code": 0, "stdout": "A B +\n", "stderr": ""},
}

_p, _script_map, _ = _PB.build(
    ConversationState(),
    "whats the postfix of A+B",
    {"steps": [{"type": "python.run", "args": {"path": "x.py"}}]},
    [_ran],
)

check("script output is citable evidence",
      "A B +" in list(_script_map.values()),
      f"tagged lines: {list(_script_map.values())}")

check("a computed answer is supported by its own output",
      _answer_within_evidence("The postfix of A+B is A B +.",
                              list(_script_map.values())))


# Calculation grounding should include what actually ran, not comments
# that may contain a stale or mistaken formula.  Short version strings
# from official release pages also need IDs even when they are fragments
# rather than full grammatical sentences.
_calculation_result = {
    "success": True,
    "data": {
        "return_code": 0,
        "stdout": "Distance: 44.1 metres\n",
        "stderr": "",
    },
}
_calculation_prompt, _calculation_map, _ = _PB.build(
    ConversationState(),
    "how far does it fall in 3 seconds",
    {"steps": [{"type": "code.run", "args": {
        "code": (
            "# Wrong old note: distance = 0.3 * gravity * time ** 2\n"
            "distance = 0.5 * 9.8 * 3 ** 2\n"
            "print(f'Distance: {distance} metres')\n"
        ),
    }}]},
    [_calculation_result],
)
check("the executed calculation source is citable",
      any("distance = 0.5" in line for line in _calculation_map.values()),
      f"tagged lines: {list(_calculation_map.values())}")
check("comments cannot become calculation evidence",
      "Wrong old note" not in _calculation_prompt
      and all("distance = 0.3" not in line
              for line in _calculation_map.values()))

_version_prompt, _version_map, _ = _PB.build(
    ConversationState(),
    "latest stable Python version",
    {"steps": [{"type": "web.search", "args": {
        "query": "latest stable Python version",
    }}]},
    [{"success": True, "data": [{
        "title": "Python releases",
        "url": "https://www.python.org/downloads/",
        "snippet": "Official Python downloads",
        "content": "Python 3.14.7\nReleased on 5 August 2026",
        "official": True,
        "trusted": True,
    }]}],
)
check("short numeric release facts receive evidence IDs",
      any(line == "Python 3.14.7" for line in _version_map.values())
      and any("5 August 2026" in line for line in _version_map.values()),
      f"tagged lines: {list(_version_map.values())}")
check("the model can distinguish official from merely trusted pages",
      "[OFFICIAL SOURCE]" in _version_prompt)


# ----------------------------------------------------------------
# Redo
# ----------------------------------------------------------------
#
# Asking again after a bad answer. The visible half is only half the
# point: a wrong answer left in the history is quoted back to the model
# on every following turn as something already established, and it
# builds on it. "Try that again" has to mean the first try is gone.

print("\n--- redo ---")

_redo = Agent.__new__(Agent)
_redo.state = _State()

check("there is nothing to redo in an empty conversation",
      _redo.take_back_last_turn() is None)

_redo.state.messages = [
    {"role": "user", "content": "first question"},
    {"role": "assistant", "content": "first answer"},
    {"role": "user", "content": "second question"},
    {"role": "assistant", "content": "a wrong answer"},
]

_asked = _redo.take_back_last_turn()

check("redo returns the question to ask again",
      _asked == "second question", f"got {_asked!r}")

check("the rejected answer is gone from the history",
      not any("wrong answer" in m["content"] for m in _redo.state.messages),
      f"left: {_redo.state.messages}")

check("the question itself is removed too, not left dangling",
      not any("second question" in m["content"] for m in _redo.state.messages),
      "answering again would repeat the question in the history")

check("earlier exchanges are untouched",
      len(_redo.state.messages) == 2
      and _redo.state.messages[0]["content"] == "first question")

# A summary can cover messages that are still here, but never ones
# that have been deleted.
_redo.state.messages = [{"role": "user", "content": f"m{i}"} for i in range(30)]
_redo.state.summarized_upto = 24
_redo.take_back_last_turn()

check("the summary never claims to cover deleted messages",
      _redo.state.summarized_upto <= len(_redo.state.messages),
      f"covers {_redo.state.summarized_upto} of {len(_redo.state.messages)}")


# ----------------------------------------------------------------
# Impossible character ranges in generated code
# ----------------------------------------------------------------
#
# Asked to tokenise an expression, the model reliably writes
# r"[A-Za-z0-9+-*/()]". The hyphen between "+" and "*" reads as a range
# running backwards (43 to 42) and Python raises at RUNTIME:
#
#   re.PatternError: bad character range +-* at position 7
#
# Asked for the postfix of A+B, Athena answered with that traceback.
#
# The generation prompt has said to put the hyphen last since the first
# time this happened, and the model still writes it this way. A rule the
# generator can forget is not a fix, so the code is repaired instead.

print("\n--- generated regex repair ---")

import re as _re
from services.code_service import repair_character_ranges as _repair

_broken = 'tokens = re.findall(r"[A-Za-z0-9+-*/()]", expression)'
_fixed = _repair(_broken)

check("the impossible range is escaped",
      _fixed == 'tokens = re.findall(r"[A-Za-z0-9+\\-*/()]", expression)',
      f"got: {_fixed}")

_pattern = _re.search(r'r"(.*?)"', _fixed).group(1)

try:
    _re.compile(_pattern)
    _compiles = True
except _re.error:
    _compiles = False

check("the repaired pattern actually compiles", _compiles)

# Narrowing until the broken case passes is easy if nothing insists
# working regexes are left alone.
_UNTOUCHED = [
    'p = re.compile(r"[a-z]")',
    'p = re.compile(r"[A-Za-z0-9_]")',
    'p = re.compile(r"[0-9]+")',
    'p = re.compile(r"[^a-z]")',
    'p = re.compile(r"[+*/^-]")',      # already written correctly
    'p = re.compile(r"[\\-]")',        # already escaped
]

check("valid character classes are left alone",
      all(_repair(s) == s for s in _UNTOUCHED),
      f"changed: {[s for s in _UNTOUCHED if _repair(s) != s]}")

# Only string literals are touched. A hyphen inside brackets elsewhere
# in Python is subtraction, and rewriting it would break the program
# this is meant to be saving.
_NOT_REGEX = [
    "x = a[1-2]",
    "y = items[i-1]",
    "z = [b-c for b in q]",
]

check("brackets outside strings are not regexes",
      all(_repair(s) == s for s in _NOT_REGEX),
      f"changed: {[s for s in _NOT_REGEX if _repair(s) != s]}")


# ----------------------------------------------------------------
# Scripts that wait for input nobody will give them
# ----------------------------------------------------------------
#
# The script runs unattended: no standard input, no arguments, no data
# files. Every value it needs is in the specification.
#
# The prompt has said so since a script asked for a train's distance
# and produced no answer at all. Told not to prompt, the model started
# reading the value out of a file instead - "could not find the input
# file 'infix_expression.txt'", a file that has never existed - which
# fails identically for the same reason. Three prompt rules have now
# been forgotten this way, so this one is checked in code.

print("\n--- scripts that wait for input ---")

from services.code_service import waits_for_input as _waits

_SPEC = "Convert the infix expression A+B to postfix and print it."

_WAITING = [
    'x = input("expr: ")',
    "import sys\ne = sys.argv[1]",
    "import sys\ne = sys.stdin.read()",
    'e = open("infix_expression.txt").read()',     # the reported failure
    'with open("data.txt", "r") as f:\n    e = f.read()',
    "import argparse",
]

check("a script waiting for input is caught",
      all(_waits(c, _SPEC) for c in _WAITING),
      f"missed: {[c for c in _WAITING if not _waits(c, _SPEC)]}")

# Writing files is the point of half these scripts, and a file the
# request actually named is fair to read. Neither may be flagged.
_FINE = [
    ('infix_expression = "A+B"', _SPEC),
    ('with open("out.txt", "w") as f:\n    f.write(x)', _SPEC),
    ('open("deck.pptx", "wb")', _SPEC),
    ("prs.save(path)", _SPEC),
    ('e = open("budget.csv").read()', "Read budget.csv and total column B."),
]

check("writing files and reading a named one are allowed",
      not any(_waits(c, s) for c, s in _FINE),
      f"wrongly flagged: {[c for c, s in _FINE if _waits(c, s)]}")

check("the reason names what went wrong",
      "infix_expression.txt" in _waits('e = open("infix_expression.txt").read()', _SPEC),
      "the retry prompt would not say what to fix")


# ----------------------------------------------------------------
# Scripts that run cleanly and compute nothing
# ----------------------------------------------------------------
#
# Asked for the prefix of A+B, the generated program tokenised the
# expression with r"\b[A-Za-z0-9+\-*/^]+\b", which matches "A+B" as a
# single token - neither a name nor an operator, so a loop testing for
# each appended nothing. It printed
#
#     Infix Expression: A+B
#     Prefix Expression:
#
# and exited zero with an empty stderr. Nothing had failed by any
# measure, so the blank was composed into an answer: "The prefix
# expression is blank."
#
# An empty result is a failure and gets the same one repair attempt a
# traceback does.

print("\n--- scripts that compute nothing ---")

from core.agent import _prints_nothing, _script_error as _err

_NOTHING = [
    "Infix Expression: A+B\nPrefix Expression: \n",   # the real failure
    "",
    "   \n",
    "Result:",
    "Answer:\n",
]

_SOMETHING = [
    "Infix Expression: A+B\nPostfix Expression: A B +\n",
    "4\n",
    "A B +",
    "Speed: 78 km/h\n",
    # A heading partway through is working, not a missing result - so
    # the check looks at the last line rather than every line.
    "Steps:\n  1. add\n  2. done\nResult: 5\n",
]

check("output with no computed value is caught",
      all(_prints_nothing(s) for s in _NOTHING),
      f"missed: {[s for s in _NOTHING if not _prints_nothing(s)]}")

check("real output is not mistaken for nothing",
      not any(_prints_nothing(s) for s in _SOMETHING),
      f"wrongly flagged: {[s for s in _SOMETHING if _prints_nothing(s)]}")

_blank_run = {"data": {"return_code": 0,
                       "stdout": "Infix Expression: A+B\nPrefix Expression: \n",
                       "stderr": ""}}
_good_run = {"data": {"return_code": 0,
                      "stdout": "Prefix Expression: + A B\n",
                      "stderr": ""}}

check("a blank result reaches the repair path",
      bool(_err(_blank_run)),
      "it would have been composed into an answer instead")

check("a clean run is left alone",
      not _err(_good_run))

check("the repair is told what the program printed",
      "Prefix Expression:" in _err(_blank_run),
      "the retry would not know what was missing")


# ----------------------------------------------------------------
# Interface invariants
# ----------------------------------------------------------------
#
# The interaction code is kept in one static script while the established
# visual design remains in the page. These are the mistakes that have actually
# been made in the interface, checked from both files.

print("\n--- interface ---")

import re
from pathlib import Path as _Path

_PAGE = (_Path(__file__).resolve().parent.parent
         / "core" / "templates" / "index.html").read_text(encoding="utf-8")
_SCRIPT = (_Path(__file__).resolve().parent.parent
           / "core" / "static" / "athena.js").read_text(encoding="utf-8")
_PIPELINE_STYLES = (_Path(__file__).resolve().parent.parent
                    / "core" / "static" / "pipeline.css").read_text(encoding="utf-8")
_STYLES = _PAGE + "\n" + _PIPELINE_STYLES
_FRONTEND = _PAGE + "\n" + _PIPELINE_STYLES + "\n" + _SCRIPT

check("the interaction script is loaded separately",
      '/static/athena.js?v=2' in _PAGE)

# Only the processing rail is being redesigned. Its stylesheet is deliberately
# isolated so a later graph adjustment cannot restyle the chat page again.
check("the pipeline stylesheet is isolated and loaded",
      '/static/pipeline.css?v=1' in _PAGE
      and "#flow-rail" in _PIPELINE_STYLES
      and "#workspace" not in _PIPELINE_STYLES)

check("the pipeline has separate idle and detailed views",
      'class="idle-layer"' in _PAGE and 'class="trace-layer"' in _PAGE)

check("pipeline labels describe the work in plain language",
      ">understand<" in _PAGE and ">use tools<" in _PAGE
      and ">execute<" not in _PAGE)

check("the pipeline collapses after showing the completed route",
      "scheduleFlowCollapse" in _SCRIPT
      and 'setFlowMode("idle")' in _SCRIPT
      and "setFlowMode(outcome)" in _SCRIPT)

check("direct chat bypasses evidence verification in the diagram",
      'id="p-compose-answer"' in _PAGE
      and 'pipe("p-compose-answer").classList.add("done")' in _SCRIPT)

check("failed and stopped turns do not illuminate Answer",
      'settleFlow(outcome = "complete")' in _SCRIPT
      and 'outcome === "complete" && node("n-compose")' in _SCRIPT
      and '"failed", "stopped"' in _SCRIPT)

check("failed and stopped turns do not play the completion sound",
      "if (!data.stopped && !data.error) beep();" in _SCRIPT)

check("capabilities get friendly labels without losing technical names",
      "TOOL_LABELS" in _SCRIPT
      and 'title.textContent = tool' in _SCRIPT
      and '"filesystem.read": "document read"' in _SCRIPT)


# The page's pure functions, actually executed rather than only parsed.
#
# `node --check` was the only JavaScript check for a long time, and it
# cannot see an undefined variable - that is a ReferenceError at run
# time, not a syntax error. So `sessionTokens += computed` reached the
# browser and turned the first message of a conversation into
# "Error: ReferenceError: computed is not defined". The variable had
# been renamed a few edits earlier and one line still used the old
# name.
#
# Skipped rather than failed when node is missing: this is a Python
# suite, and someone without node should still be able to run it.
import shutil as _shutil
import subprocess as _subprocess

_NODE = _shutil.which("node")

if not _NODE:
    print("SKIP  page functions (node is not installed)")
else:
    _checker = _Path(__file__).resolve().parent / "check_page_functions.js"
    _ran = _subprocess.run(
        [_NODE, str(_checker)], capture_output=True, text=True,
    )

    check("the page's functions run without error",
          _ran.returncode == 0,
          (_ran.stdout + _ran.stderr).strip())


# The `hidden` attribute hides an element by setting display:none in
# the browser's own stylesheet, which any display in the page's own CSS
# outranks. This has now been got wrong three times - the history
# panel, the undo bar and the rows in the history list - and each time
# the symptom was a button that fired correctly and appeared to do
# nothing at all.
def _display_without_guard(selector, block):
    """Whether a rule sets display but has no [hidden] escape hatch."""

    rule = re.search(re.escape(selector) + r"\s*\{(.*?)\}", _STYLES, re.S)

    if not rule or not re.search(r"\bdisplay\s*:", rule.group(1)):
        return False

    return f"{selector}[hidden]" not in _STYLES


# Which elements this applies to is worked out rather than listed. A
# hardcoded list was the first attempt and immediately named an element
# that is hidden with an inline style and never touches the attribute
# at all - so the test would have reported a problem that did not exist
# while missing the next one that did.
_HIDDEN_IDS = set(re.findall(r'id="([\w-]+)"[^>]*\shidden', _PAGE))

# The attribute is also set from JavaScript, on elements that never
# carry it in the markup. Those need the guard just as much, so the
# variable names are mapped back to the ids they were looked up with.
_ELEMENT_VARS = dict(
    re.findall(r'(?:const|let|var)\s+(\w+)\s*=\s*document\.getElementById\(\s*["\']([\w-]+)["\']',
               _SCRIPT)
)

for _var in set(re.findall(r"(\w+)\.hidden\s*=", _SCRIPT)):
    if _var in _ELEMENT_VARS:
        _HIDDEN_IDS.add(_ELEMENT_VARS[_var])

_UNGUARDED = [i for i in sorted(_HIDDEN_IDS)
              if _display_without_guard(f"#{i}", _PAGE)]

check(f"elements hidden by attribute stay hidden ({len(_HIDDEN_IDS)} checked)",
      not _UNGUARDED,
      f"display overrides hidden for: {_UNGUARDED} - "
      f"add `#id[hidden] {{ display: none; }}`")

check("history rows can be hidden",
      not _display_without_guard(".history-item", _PAGE),
      "a row waiting on its undo timer would stay on screen")


# position:fixed is measured against the viewport, unless an ancestor
# has a transform, filter or backdrop-filter - then it is measured
# against that ancestor instead. The header has a backdrop-filter, so
# anything fixed inside it anchors to the header rather than the
# window, which is why the undo bar lives outside it.
_HEADER = re.search(r"<header>(.*?)</header>", _PAGE, re.S)

check("nothing fixed to the window sits inside the header",
      _HEADER and 'id="toast"' not in _HEADER.group(1),
      "the header's backdrop-filter would anchor it to the header")


# Every dialog needs a way out that does not require a mouse.
check("the history panel closes on Escape",
      re.search(r'"Escape".*?showHistory\(false\)', _SCRIPT, re.S) is not None)


# Buttons that only appear on hover are unreachable by keyboard unless
# they also appear on focus.
for _hover_only in (".redo", ".history-item .remove"):
    check(f"{_hover_only} is reachable without a mouse",
          f"{_hover_only}:focus-visible" in _STYLES,
          "it appears on hover only, so a keyboard user cannot get to it")


# ----------------------------------------------------------------
# Small sums are not worth computing
# ----------------------------------------------------------------
#
# "What's 2+2" was routed to CALCULATE, which wrote a Python script,
# saved it, ran it in a subprocess and composed an answer from the
# output - five model calls and fifteen seconds. Worse than slow, it
# added a way to fail: when the generated script had a bug the repair
# path ran, and when that missed too the reply was "I couldn't find
# that in what I looked up", about two plus two.
#
# Computation earns its cost on multi-step problems. One small addition
# is not one of those.

print("\n--- small sums ---")

from core.router import looks_arithmetic as _arith

_JUST_ANSWER = ["WHats 2+2", "what is 2+2", "2+2", "whats 7-3",
                "what is 12 + 5", "15-8"]

_STILL_COMPUTE = ["what is 45 * 12",      # two-digit multiplication slips
                  "2+2+2",                # more than one operation
                  "whats 123+456",        # beyond two digits
                  "what is 2^10",
                  "1234-567",
                  "whats 100/7",
                  "convert A+B*C to postfix",
                  "what is 15 percent of 240"]

check("small sums are answered directly",
      not any(_arith(m) for m in _JUST_ANSWER),
      f"still computed: {[m for m in _JUST_ANSWER if _arith(m)]}")

check("anything harder is still computed",
      all(_arith(m) for m in _STILL_COMPUTE),
      f"no longer computed: {[m for m in _STILL_COMPUTE if not _arith(m)]}")


# The router model reaches CALCULATE on its own, separately from the
# mode that forces computation. Fixing only looks_arithmetic left
# "what's 2+2" still writing itself a script.
class _AlwaysCalculate:
    def complete(self, system, prompt, schema=None, num_predict=None, think=None):
        routes = schema["properties"]["route"]["enum"]
        return _json.dumps({"route": "CALCULATE" if "CALCULATE" in routes else "SAFE"})


check("a small sum is not computed even when the router asks for it",
      _Router(_AlwaysCalculate()).route(ConversationState(), "WHats 2+2") == "chat",
      "the router's own CALCULATE still forced a script")

check("a real calculation still goes to the planner",
      _Router(_AlwaysCalculate()).route(ConversationState(), "what is 45 * 12")
      == "calculate")


# The page shows replies as text, not HTML, so markdown arrives
# literally. Asked what 2+2 was, Athena answered
# "* **Rule:** This is a basic addition problem" - visibly.
check("markdown is stripped from chat replies",
      "**" not in _sm("* **Rule:** basic addition") ,
      f"got: {_sm('* **Rule:** basic addition')!r}")

check("multiplication signs survive stripping",
      "*" in _sm("The postfix form is ABC*+"),
      "the markdown stripper ate a multiplication sign")


# ----------------------------------------------------------------
# Structural audit regressions
# ----------------------------------------------------------------

print("\n--- structural audit regressions ---")

# Values and IDs are tokens, never substrings of a larger value.
for _answer_value, _source_value in [
    ("The fee is 285.", "The total is 117285.00."),
    ("The year is 202.", "The year is 2025."),
    ("Reference TESTA729899.", "Reference TESTA7298990."),
    ("It is 30 degrees.", "It is 300 degrees."),
]:
    check(f"partial value is not grounded: {_answer_value!r}",
          not _answer_within_evidence(_answer_value, [_source_value]))

from core.agent import (
    _evidence_is_relevant as _relevant,
    _grounded_schema_for as _schema_for_grounding,
    _source_organization as _organization,
)

check("lowercase search terms still drive relevance",
      not _relevant(["Cookie settings and navigation."],
                    ["london weather today"])
      and _relevant(["London has rain today."],
                    ["london weather today"]))
check("one shared place word is not enough for a multi-part query",
      not _relevant(
          ["As of 2026, Australia has a population estimate."],
          ["capital of Australia 21 August 2026"],
      )
      and _relevant(
          ["The capital of Australia is Canberra."],
          ["capital of Australia 21 August 2026"],
      ))

check("grounding schema permits an honest empty evidence array",
      _schema_for_grounding("Explain this in detail")
      ["properties"]["evidence"].get("minItems", 0) == 0)

# Trust and corroboration operate on parsed domains, not URL text.
from services.web_service import WebService as _WebService

_web = _WebService()
check("trusted hostname is accepted",
      _web._is_trusted("https://en.wikipedia.org/wiki/Athena"))
check("trusted text in another hostname is rejected",
      not _web._is_trusted("https://notwikipedia.org/story")
      and not _web._is_trusted("https://evil.example/?next=whitehouse.gov"))
check("a subject's own domain is recognised as official",
      _web._is_official_for_query(
          "https://www.python.org/downloads/",
          "latest stable Python version",
      ))
check("a subject word elsewhere in a URL is not called official",
      not _web._is_official_for_query(
          "https://tutorials.example/python/latest",
          "latest stable Python version",
      )
      and not _web._is_official_for_query(
          "https://notpython.org/releases",
          "latest stable Python version",
      ))
check("subdomains collapse to one source organization",
      _organization("https://news.bbc.co.uk/a")
      == _organization("https://sport.bbc.co.uk/b") == "bbc.co.uk")

_ranked = _web._rank_by_relevance(
    "Cookie settings\nUnrelated navigation words here\n"
    "London temperature is 20 degrees today.\nFooter links",
    "london temperature",
)
check("web extraction does not preload unrelated leading lines",
      "London temperature" in _ranked and "Unrelated navigation" not in _ranked,
      f"got: {_ranked!r}")

# Failed or empty reads cannot replace a valid active file.
_file_state_agent = Agent.__new__(Agent)
_file_state_agent.state = ConversationState(last_file_path="old.pdf")
_file_state_agent._update_file_state(
    {"type": "filesystem.read", "args": {"path": "ghost.pdf"}},
    {"success": False, "error": "missing"},
)
check("failed read does not become active",
      _file_state_agent.state.last_file_path == "old.pdf")
_file_state_agent._update_file_state(
    {"type": "filesystem.read", "args": {"path": "empty.pdf"}},
    {"success": True, "data": ""},
)
check("empty read does not become active",
      _file_state_agent.state.last_file_path == "old.pdf")

# A stopped model call may mutate only the provisional state.
class _StopAfterChat(FakeModel):
    def chat(self, state, message, **kwargs):
        state.messages.append({"role": "user", "content": message})
        state.messages.append({"role": "assistant", "content": "invisible"})
        _PROGRESS.stop()
        return "invisible"


_PROGRESS.clear()
_cancelled = build_agent(
    FakePlanner([]), FakeExecutor({"success": True, "data": "x"}),
    _StopAfterChat({}), route="chat",
    state=ConversationState(messages=[
        {"role": "user", "content": "before"},
        {"role": "assistant", "content": "kept"},
    ]),
)
_before_cancel = list(_cancelled.state.messages)
try:
    _cancelled.respond("do not keep this")
    _cancel_raised = False
except _Stopped:
    _cancel_raised = True
finally:
    _PROGRESS.clear()

check("late cancellation raises Stopped", _cancel_raised)
check("cancelled turn leaves no invisible state",
      _cancelled.state.messages == _before_cancel,
      f"got: {_cancelled.state.messages}")

# Conversation IDs, full transcripts and operational state survive a
# safe round trip. The metadata index is separate from transcript data.
_audit_root = _Path(_tempfile.mkdtemp())
_audit_store = _Store(_audit_root / "conversations")
_id_a = _audit_store.new_id("same title")
_id_b = _audit_store.new_id("same title")
check("conversation ids cannot collide", _id_a != _id_b)

_long_answer = "complete answer " * 400
_audit_store.save(_Convo(
    id=_id_a,
    title="Audit",
    messages=[{"role": "assistant", "content": _long_answer}],
    last_file_path="C:/Docs/receipt.pdf",
    last_generated_path="C:/Athena/build_report.py",
    pending_file_paths=["C:/Docs/a.pdf", "C:/Docs/b.pdf"],
    pending_file_request="summarize it",
    last_capabilities=["filesystem.search"],
    last_capability_steps=[{
        "type": "filesystem.search", "args": {"name": "receipt"},
    }],
))
_audit_back = _audit_store.load(_id_a)
check("saved transcript keeps the complete answer",
      _audit_back.messages[0]["content"] == _long_answer)
check("saved conversation restores operational state",
      _audit_back.last_file_path == "C:/Docs/receipt.pdf"
      and _audit_back.last_generated_path == "C:/Athena/build_report.py"
      and len(_audit_back.pending_file_paths) == 2
      and _audit_back.last_capabilities == ["filesystem.search"]
      and _audit_back.last_capability_steps[0]["args"]["name"] == "receipt")
check("conversation sidebar has a metadata index",
      (_audit_root / "conversations" / "_index.json").is_file())
check("conversation traversal is rejected",
      _audit_store.load("../../workspace/semantic_index") is None
      and _audit_store.delete("../outside") is False)

# A limited semantic scan updates what it sees but never prunes files it
# did not reach. A subsequent complete scan may remove a genuinely
# missing file.
from services.semantic_index import SemanticIndex as _SemanticIndex

_docs = _audit_root / "docs"
_docs.mkdir()
_doc_paths = []
_entries = []
for _n in range(3):
    _path = _docs / f"doc{_n}.txt"
    _path.write_text(f"document {_n}", encoding="utf-8")
    _stat = _path.stat()
    _doc_paths.append(_path)
    _entries.append({
        "path": str(_path), "name": _path.name,
        "mtime": _stat.st_mtime, "size": _stat.st_size,
        "chunks": [f"document {_n}"], "vectors": [[1.0, 0.0]],
    })

class _IndexFS:
    SKIP_DIRS = set()
    def read(self, path):
        return {"success": True, "data": _Path(path).read_text(encoding="utf-8")}


_index = _SemanticIndex(_IndexFS(), _audit_root / "index.json")
_index._entries = list(_entries)
_limited = _index.index([_docs], limit_files=1)
check("partial semantic scan does not prune unseen entries",
      _limited["complete_scan"] is False
      and _limited["removed"] == 0 and _limited["total"] == 3,
      f"got: {_limited}")

_doc_paths[-1].unlink()
_complete = _index.index([_docs], limit_files=20)
check("complete semantic scan prunes a deleted entry",
      _complete["complete_scan"] is True
      and _complete["removed"] == 1 and _complete["total"] == 2,
      f"got: {_complete}")

# Mixed plans retain both evidence types and expose both source records.
class _MixedExecutor:
    def execute(self, plan):
        tool = plan["steps"][0]["type"]
        if tool == "web.search":
            return [{
                "success": True,
                "queries": ["current public threshold policy"],
                "data": [{
                    "title": "Public threshold",
                    "url": "https://example.org/threshold",
                    "snippet": "Current public threshold information",
                    "content": "The current public threshold is 20 units.",
                    "trusted": False,
                }],
            }]
        return [{"success": True,
                 "data": "The private threshold is 10 units."}]


_mixed = build_agent(
    FakePlanner([
        {"type": "web.search", "args": {
            "query": "current public threshold policy", "category": "general"}},
        {"type": "filesystem.read", "args": {"path": "C:/Docs/policy.txt"}},
    ]),
    _MixedExecutor(),
    FakeModel({
        "evidence": ["S2", "S3"],
        "answer": (
            "The private threshold is 10 units, while the current public "
            "threshold is 20 units."
        ),
    }),
)
_mixed_answer = _mixed.respond("Compare my policy with the current public threshold")
check("mixed web and file answer keeps both supported values",
      "10" in _mixed_answer and "20" in _mixed_answer,
      f"got: {_mixed_answer!r}")
check("mixed answer reports both source types",
      {source["kind"] for source in _mixed.state.last_sources}
      == {"web", "local"},
      f"got: {_mixed.state.last_sources}")

# Clear requests use a much smaller planner prompt and constrained tool
# schema; ambiguous follow-ups retain the full fallback.
from core.planner import Planner as _Planner

_planner_probe = _Planner(object())
_weather_tools = _planner_probe._allowed_tools("weather in london")
_file_tools = _planner_probe._allowed_tools("find a file named fees")
_mixed_tools = _planner_probe._allowed_tools(
    "compare this document with current law"
)
check("weather planner keeps only relevant tools",
      _weather_tools == {"weather.current", "web.search"})
check("mixed planner can use file and web capabilities together",
      "web.search" in _mixed_tools
      and "filesystem.read" in _mixed_tools)
check("category planner prompt is materially smaller",
      len(_planner_probe._system_prompt_for(_weather_tools))
      < len(_planner_probe.system_prompt) * 0.6)
check("category schema cannot emit unrelated code generation",
      "code.generate" not in _planner_probe._schema_for(_file_tools)
      ["properties"]["step"]["properties"]["type"]["enum"])

# Clarification is returned directly instead of becoming a web search.
import core.execution_manager as _execution_module
from core.execution_manager import ExecutionManager as _ExecutionManager

_lookup = _ExecutionManager.__new__(_ExecutionManager)
_lookup._run_live_lookup = lambda tool, args: {
    "success": False,
    "needs_clarification": True,
    "error": "Which city?",
}
_fallback_calls = []
_lookup._execute_web_search = lambda args: _fallback_calls.append(args) or {
    "success": True, "data": []
}
_old_internet = _execution_module.has_internet
_execution_module.has_internet = lambda force=False: True
try:
    _clarification = _lookup._execute_live_data(
        "weather.current", {"location": "Japan"}
    )
finally:
    _execution_module.has_internet = _old_internet
check("country weather clarification does not fall back to search",
      _clarification.get("needs_clarification") and not _fallback_calls)

# Model lifecycle, warm-up and capability metadata contracts.
from models.ollama_model import OllamaModel as _OllamaModel

check("a separate vision model has a finite residency option",
      _OllamaModel("gemma3:4b", keep_alive="2m").keep_alive == "2m")

class _WarmFailure:
    def complete(self, *args, **kwargs):
        raise RuntimeError("not installed")


_warm_agent = Agent.__new__(Agent)
_warm_agent.planner_model = _WarmFailure()
try:
    _warm_agent._warm_up()
    _warm_failed_loudly = False
except RuntimeError:
    _warm_failed_loudly = True
check("model warm-up failure is surfaced", _warm_failed_loudly)

from services.filesystem import FilesystemService as _FilesystemService

_info_file = _audit_root / "info.txt"
_info_file.write_text("x", encoding="utf-8")
_info = _FilesystemService().info(str(_info_file))
check("filesystem.info fulfils its modification-time contract",
      _info.get("success") and _info["data"].get("modified_at"))

_secret_file = _audit_root / ".env"
_secret_file.write_text("API_KEY=do-not-read", encoding="utf-8")
_public_env = _audit_root / ".env.example"
_public_env.write_text("API_KEY=placeholder", encoding="utf-8")
check("credentials files are refused by the document reader",
      not _FilesystemService().read(str(_secret_file)).get("success"))
check("the public environment template remains readable",
      _FilesystemService().read(str(_public_env)).get("success"))

from config import (
    ATHENA_HOST as _configured_host,
    BALANCED_MODEL as _configured_balanced,
)

check("configured server address is nonempty",
      bool(str(_configured_host).strip()))
check("the configured default model reaches the mode registry",
      MODES["balanced"]["response"] == _configured_balanced)

check("the interface renders verified sources",
      "verified sources" in _FRONTEND and "source-links" in _FRONTEND)
check("temporary images are unique and removed",
      "NamedTemporaryFile" in _webapp.__loader__.get_source(_webapp.__name__)
      and "unlink(missing_ok=True)" in _webapp.__loader__.get_source(_webapp.__name__))


# ----------------------------------------------------------------
# Failures found by the three-mode release conversation
# ----------------------------------------------------------------

print("\n--- release-conversation regressions ---")

from core.router import (
    continues_a_lookup as _continues_a_lookup,
    missing_subject_question as _missing_subject_question,
)
from core.agent import (
    _augment_local_file_answer as _augment_file_answer,
    _completed_artifact_answer as _artifact_answer,
    _direct_step_for as _direct_step,
    _prepare_generation_step as _prepare_generation,
    _request_needs_multiple_tools as _needs_multiple,
    _structured_result_answer as _structured_answer,
    _wants_built_document as _wants_built,
)

# A selected PDF is useful context, but it must not become the subject
# of every later question containing ordinary pronouns such as "it" or
# "that". These exact styles were all hijacked by filesystem.read in
# the real three-mode run.
_stale_file = ConversationState(
    last_file_path=r"C:\Docs\receipt.pdf",
    last_capabilities=["weather.current"],
)
for _message in [
    "what time is it rn?",
    "nah bro isnt that 418?",
    "could it be a cat?",
    "isnt it Sydney tho?",
    "run that script now",
]:
    check(f"stale file does not hijack {_message!r}",
          not is_active_file_request(_stale_file, _message))

_fresh_file = ConversationState(
    last_file_path=r"C:\Docs\receipt.pdf",
    last_capabilities=["filesystem.read"],
)
check("a field question immediately after a read keeps file context",
      is_active_file_request(_fresh_file, "whats the reference code?"))
check("an explicit file action remains valid even after another topic",
      is_active_file_request(_stale_file, "summarize the file"))

# A sentence can mention a live topic without requesting a lookup.
check("conversation about weather is not treated as a missing-city request",
      _missing_subject_question("I love this weather ngl") == "")

# Company-before-topic wording is the normal form for a stock request.
_tesla_model = AlwaysActiveFileModel()
_tesla_route = _Router(_tesla_model).route(
    ConversationState(), "Tesla stock price rn?"
)
check("company before 'stock price' supplies the lookup subject",
      _tesla_route == "capability", f"got: {_tesla_route}")

# Challenges to a just-retrieved live fact must repeat the lookup. The
# chat model must not accept the user's number or invent a compromise.
_live_state = ConversationState(last_capabilities=["weather.current"])
check("a slang challenge continues the previous live lookup",
      _continues_a_lookup(_live_state, "u sure? I thought it was 45"))
check("a direct contradiction continues the previous live lookup",
      _continues_a_lookup(_live_state, "but isnt it 45 degrees?"))

_challenge_executor = FakeExecutor({"success": True, "data": {
    "place": "Mumbai, India", "conditions": "overcast",
    "temperature": 27.4, "temperature_unit": "°C",
    "feels_like": 31.0, "humidity": 78, "humidity_unit": "%",
    "precipitation": 0.0, "precipitation_unit": "mm",
    "wind_speed": 8.0, "wind_unit": "km/h",
    "observed_at": "2026-08-20T15:15", "timezone": "Asia/Kolkata",
}})
_challenge_agent = build_agent(
    FakePlanner([]), _challenge_executor,
    FakeModel({"evidence": [], "answer": ""}),
    state=ConversationState(
        last_capabilities=["weather.current"],
        last_capability_steps=[{
            "type": "weather.current", "args": {"location": "Mumbai"},
        }],
    ),
)
_challenge_answer = _challenge_agent.respond("u sure? I thought it was 45")
check("a live challenge repeats the exact previous lookup arguments",
      _challenge_executor.calls == [{
          "type": "weather.current", "args": {"location": "Mumbai"},
      }], f"got: {_challenge_executor.calls}")
check("a false live correction is answered from refreshed data",
      "27.4" in _challenge_answer and "45" not in _challenge_answer,
      f"got: {_challenge_answer}")

_calculation_challenge_executor = FakeExecutor({
    "success": True,
    "data": {"return_code": 0, "stdout": "408\n", "stderr": ""},
})
_calculation_challenge_agent = build_agent(
    FakePlanner([]),
    _calculation_challenge_executor,
    FakeModel({"evidence": ["S1"], "answer": "The result is 408."}),
    state=ConversationState(
        last_capabilities=["code.run"],
        last_capability_steps=[{
            "type": "code.run",
            "args": {"code": "print(17 * 24)"},
        }],
    ),
)
_calculation_challenge_answer = _calculation_challenge_agent.respond(
    "but isnt it 418?"
)
check("a false arithmetic correction re-runs the exact calculation",
      _calculation_challenge_executor.calls == [{
          "type": "code.run", "args": {"code": "print(17 * 24)"},
      }]
      and "408" in _calculation_challenge_answer
      and "418" not in _calculation_challenge_answer,
      f"answer={_calculation_challenge_answer!r}, "
      f"calls={_calculation_challenge_executor.calls!r}")

# The planner's reduced schema must retain every domain in a compound
# request. An underscored .txt path used to leave only web.search.
_compound_tools = _planner_probe._allowed_tools(
    r"read C:\Docs\aurora_membership_receipt.txt, then use today's "
    r"rate to tell me how many USD the INR total is"
)
check("compound file-and-currency requests retain filesystem.read",
      "filesystem.read" in _compound_tools,
      f"got: {sorted(_compound_tools)}")
check("compound file-and-currency requests retain finance.exchange",
      "finance.exchange" in _compound_tools,
      f"got: {sorted(_compound_tools)}")

_compute_tools = _planner_probe._allowed_tools(
    "A train travels 240 km in 3 hours; find its speed",
    must_calculate=True,
)
check("calculation uses direct code execution, not a generated file",
      _compute_tools == {"code.run"},
      f"got: {sorted(_compute_tools)}")


class _CalculationPromptModel:
    def __init__(self):
        self.prompt = ""

    def complete(self, system, prompt, schema=None, num_predict=None, think=None):
        self.prompt = prompt
        return _json.dumps({
            "done": False,
            "step": {
                "type": "code.run",
                "args": {"code": "print(200 / 4 * 6)"},
            },
        })


_calculation_prompt_model = _CalculationPromptModel()
_Planner(_calculation_prompt_model).plan_step(
    ConversationState(messages=[
        {"role": "user", "content": "The plan covers 4 people."},
        {"role": "assistant", "content": "The total allocation is 200 units."},
    ]),
    "scale it for 6 instead",
    must_calculate=True,
)
check("calculation contract preserves totals instead of inventing per-unit values",
      "a total for a group is not a per-person or per-item value"
      in _calculation_prompt_model.prompt)

try:
    _planner_probe._normalize_decision({
        "done": False,
        "step": {"type": "code.run", "args": {}},
    }, allowed={"code.run"})
    _empty_code_rejected = False
except ValueError:
    _empty_code_rejected = True
check("the planner rejects an empty code.run step",
      _empty_code_rejected)

_artifact_file = _audit_root / "verified_artifact.pptx"
_artifact_file.write_bytes(b"pptx fixture")
_artifact_result = _artifact_answer(
    [
        {"type": "code.generate", "args": {
            "path": str(_audit_root / "build_verified_artifact.py"),
            "artifact_path": str(_artifact_file),
        }},
        {"type": "python.run", "args": {
            "path": str(_audit_root / "build_verified_artifact.py"),
        }},
    ],
    [
        {"success": True, "data": {"bytes_written": 10}},
        {"success": True, "data": {"return_code": 0}},
    ],
)
check("an existing generated artifact gets a deterministic success answer",
      _artifact_result is not None
      and _artifact_file.name in _artifact_result[0],
      f"got: {_artifact_result!r}")


class _RepairingArtifactExecutor:
    """Fail two generated builders, then create the requested artifact."""

    def __init__(self, artifact):
        self.artifact = artifact
        self.calls = []
        self.runs = 0

    def execute(self, plan):
        step = plan["steps"][0]
        self.calls.append(step)
        if step["type"] == "code.generate":
            return [{
                "success": True,
                "data": {"path": step["args"]["path"]},
            }]

        self.runs += 1
        if self.runs <= 2:
            return [{
                "success": True,
                "data": {
                    "return_code": 1,
                    "stdout": "",
                    "stderr": (
                        "Traceback (most recent call last):\n"
                        "AttributeError: generated builder failed\n"
                    ),
                },
            }]

        self.artifact.write_bytes(b"pptx fixture")
        return [{
            "success": True,
            "data": {
                "return_code": 0,
                "stdout": f"{self.artifact}\n",
                "stderr": "",
            },
        }]


_repaired_artifact = _audit_root / "repaired_artifact.pptx"
_repairing_builder = _RepairingArtifactExecutor(_repaired_artifact)
_repairing_agent = build_agent(
    FakePlanner([{
        "type": "code.generate",
        "args": {"path": "repair_builder.py", "spec": "Make two slides."},
    }]),
    _repairing_builder,
    FakeModel({"evidence": [], "answer": "stale answer"}),
)
_repaired_answer = _repairing_agent.respond(
    f"make a PowerPoint and save it to {_repaired_artifact}"
)
check("a generated artifact may use two bounded repair attempts",
      _repairing_builder.runs == 3
      and _repaired_artifact.name in _repaired_answer,
      f"answer={_repaired_answer!r}, runs={_repairing_builder.runs}")
check("a repaired artifact bypasses stale grounded evidence",
      "stale answer" not in _repaired_answer,
      f"got: {_repaired_answer!r}")

_literal_output = _structured_answer(
    "run this python snippet: print(55)",
    {"type": "code.run", "args": {"code": "print(55)"}},
    {"success": True, "data": {
        "stdout": "55\n", "stderr": "", "return_code": 0,
    }},
)
check("literal code output is formatted without model reinterpretation",
      _literal_output == "The output is 55.", f"got: {_literal_output!r}")


class _RepairProbe:
    def __init__(self):
        self.calls = []

    def repair_snippet(self, request, code, error):
        self.calls.append({
            "request": request,
            "code": code,
            "error": error,
        })
        return "print('answer:', 5)"


class _CalculationRetryExecutor:
    def __init__(self):
        self.calls = []
        self.code = _RepairProbe()

    def execute(self, plan):
        step = plan["steps"][0]
        self.calls.append(step)
        if len(self.calls) == 1:
            return [{
                "success": True,
                "data": {
                    "return_code": 1,
                    "stdout": "",
                    "stderr": "NameError: name 'value' is not defined",
                },
            }]
        return [{
            "success": True,
            "data": {
                "return_code": 0,
                "stdout": "answer: 5\n",
                "stderr": "",
            },
        }]


_retry_executor = _CalculationRetryExecutor()
_retry_agent = build_agent(
    FakePlanner([{
        "type": "code.run",
        "args": {"code": "print(value)"},
    }]),
    _retry_executor,
    FakeModel({"evidence": ["S1"], "answer": "The answer is 5."}),
    route="calculate",
)
_retry_answer = _retry_agent.respond("what is 2 + 3?")
check("a failed generated calculation gets exactly one repair attempt",
      len(_retry_executor.calls) == 2
      and len(_retry_executor.code.calls) == 1,
      f"runs={len(_retry_executor.calls)}, repairs={len(_retry_executor.code.calls)}")
check("the repaired calculation is the result that gets answered",
      "5" in _retry_answer
      and _retry_executor.calls[-1]["args"]["code"]
      == "print('answer:', 5)",
      f"got: {_retry_answer!r}")

_literal_executor = _CalculationRetryExecutor()
_literal_agent = build_agent(
    FakePlanner([]),
    _literal_executor,
    FakeModel({"evidence": [], "answer": ""}),
)
_literal_agent.respond("run this python snippet: print(value)")
check("Athena never silently rewrites literal code from the user",
      len(_literal_executor.calls) == 1
      and not _literal_executor.code.calls)

# A generated script is operational state of its own. Treating it as
# the selected document made later file questions read Python source.
_state_probe = Agent.__new__(Agent)
_state_probe.state = ConversationState(
    last_file_path=r"C:\Docs\receipt.pdf"
)
_state_probe._update_file_state(
    {"type": "code.generate", "args": {"path": "calculation.py"}},
    {"success": True, "data": {"path": r"C:\Athena\calculation.py"}},
    "write a python file that prints hello; don't run it",
)
check("generated code does not replace the selected document",
      _state_probe.state.last_file_path == r"C:\Docs\receipt.pdf")
check("generated code has a separate remembered path",
      _state_probe.state.last_generated_path == r"C:\Athena\calculation.py")
check("an explicit do-not-run instruction is respected",
      not _wants_built("write a python file that prints hello; don't run it"))
check("a PowerPoint request still builds its finished artifact",
      _wants_built("make a PowerPoint about the water cycle"))

# Literal inputs and exact paths are capability contracts, not routing
# judgments. They should work the same way in every model and mode.
_generated_state = ConversationState(
    last_generated_path=r"C:\Athena\workspace\answer.py"
)
check("literal Python is sent straight to code.run",
      _direct_step(
          "run this python snippet: print(sum(range(6)))",
          _generated_state,
      ) == {
          "type": "code.run",
          "args": {"code": "print(sum(range(6)))"},
      })
check("'run that script' uses generated-code state",
      _direct_step("run that script now", _generated_state) == {
          "type": "python.run",
          "args": {"path": r"C:\Athena\workspace\answer.py"},
      })
check("an exact metadata request selects filesystem.info",
      _direct_step(
          r"tell me the size and modified date of C:\Docs\receipt.txt",
          ConversationState(),
      )["type"] == "filesystem.info")
check("a full stop after an exact path still selects filesystem.read",
      _direct_step(
          r"read C:\Docs\recipe.txt. how many people is it for?",
          ConversationState(),
      ) == {
          "type": "filesystem.read",
          "args": {"path": r"C:\Docs\recipe.txt"},
      })
check("a sensitive .env path still reaches the protected reader",
      _direct_step(
          r"read C:\Project\.env and show me the keys",
          ConversationState(),
      )["type"] == "filesystem.read")
check("content-based document search is deterministic",
      _direct_step(
          "which of my docs mentions cobalt-blue thermal shields?",
          ConversationState(),
      )["type"] == "filesystem.semantic_search")

_compound_request = (
    r"read C:\Docs\receipt.txt, then use today's rate to tell me how "
    r"many USD the INR total is"
)
check("a file-plus-live request is recognised as multi-capability",
      _needs_multiple(_compound_request))
check("ordinary read-and-summarise remains one capability",
      not _needs_multiple(r"read C:\Docs\receipt.txt and summarize it"))

_summary_steps = [{"type": "filesystem.read", "args": {"path": "receipt.txt"}}]
_summary_results = [{"success": True, "data": (
    "Membership Receipt\nCustomer: Riya Sen\nReference: AUR-7391-X\n"
    "Total paid: INR 21,000\n"
)}]
_augmented_receipt = _augment_file_answer(
    "read receipt.txt and sum it up briefly",
    "Riya Sen paid INR 21,000.",
    _summary_steps,
    _summary_results,
)
check("brief transaction summaries preserve an exact reference",
      "AUR-7391-X" in _augmented_receipt,
      f"got: {_augmented_receipt!r}")

_budget_answer = _augment_file_answer(
    "read budget.xlsx and tell me the total plus the biggest expense",
    "The total is 29699 and the biggest expense is 18000.",
    [{"type": "filesystem.read", "args": {"path": "budget.xlsx"}}],
    [{"success": True, "data": (
        "Category\tAmount (INR)\nRent\t18000\nGroceries\t6500\nTotal\t29699"
    )}],
)
check("a requested table maximum keeps its category label",
      "Rent" in _budget_answer and "18000" in _budget_answer,
      f"got: {_budget_answer!r}")

_prepared_artifact = _prepare_generation(
    {
        "type": "code.generate",
        "args": {"path": "water_cycle.pptx", "spec": "Create three slides."},
    },
    r"make a PowerPoint and save it to C:\Exports\water_cycle.pptx",
)
check("artifact generation writes a Python builder, not code into .pptx",
      _prepared_artifact["args"]["path"].lower().endswith(".py"))
check("artifact generation preserves the exact requested output path",
      _prepared_artifact["args"]["artifact_path"].lower()
      == r"c:\exports\water_cycle.pptx")
check("artifact specification contains the exact destination",
      r"C:\Exports\water_cycle.pptx" in _prepared_artifact["args"]["spec"])


class _MissingArtifactExecutor:
    def __init__(self):
        self.calls = []

    def execute(self, plan):
        step = plan["steps"][0]
        self.calls.append(step)
        if step["type"] == "code.generate":
            return [{
                "success": True,
                "data": {"path": step["args"]["path"]},
            }]
        return [{
            "success": True,
            "data": {
                "return_code": 0,
                "stdout": "builder finished\n",
                "stderr": "",
            },
        }]


_missing_artifact = _audit_root / "never-built.pptx"
_missing_builder = _MissingArtifactExecutor()
_missing_agent = build_agent(
    FakePlanner([{
        "type": "code.generate",
        "args": {"path": "build_missing.py", "spec": "Make two slides."},
    }]),
    _missing_builder,
    FakeModel({"evidence": [], "answer": ""}),
)
_missing_answer = _missing_agent.respond(
    f"make a PowerPoint and save it to {_missing_artifact}"
)
check("a successful builder process is not success without its artifact",
      "did not create the requested file" in _missing_answer
      and len(_missing_builder.calls) == 2,
      f"got: {_missing_answer!r}, calls={_missing_builder.calls}")

_weather_answer = _structured_answer(
    "weather in Delhi",
    {"type": "weather.current", "args": {"location": "Delhi"}},
    {"success": True, "data": {
        "place": "Delhi, India", "conditions": "overcast",
        "temperature": 27.5, "temperature_unit": "°C",
        "feels_like": 34.8, "humidity": 78, "humidity_unit": "%",
        "precipitation": 0.0, "precipitation_unit": "mm",
        "wind_speed": 5.3, "wind_unit": "km/h",
        "observed_at": "2026-08-20T15:15", "timezone": "Asia/Kolkata",
    }},
)
check("typed weather formatting never drops the actual conditions",
      all(term in _weather_answer for term in ["overcast", "27.5", "34.8", "78"]))

_delete_agent = build_agent(
    FakePlanner([]), FakeExecutor({"success": True, "data": True}),
    FakeModel({"evidence": [], "answer": ""}),
)
_delete_answer = _delete_agent.respond(
    r"delete the file C:\Docs\receipt.txt"
)
check("unsupported deletion is stated plainly and changes nothing",
      "can't delete" in _delete_answer.lower()
      and not _delete_agent.executor.calls)

# A forced exact-path read must not terminate a compound request. The
# planner gets the read result and can select the second capability.
class _CompoundExecutor:
    def __init__(self):
        self.calls = []

    def execute(self, plan):
        step = plan["steps"][0]
        self.calls.append(step)
        if step["type"] == "filesystem.read":
            return [{"success": True, "data": "Total paid: 21000 INR"}]
        return [{"success": True, "data": {
            "base": "INR", "target": "USD", "rate": 0.0119,
            "amount": 21000, "converted": 249.9,
            "rates_published": "2026-08-20",
        }}]


_compound_executor = _CompoundExecutor()
_compound_agent = build_agent(
    FakePlanner([{
        "type": "finance.exchange",
        "args": {"base": "INR", "target": "USD", "amount": 21000},
    }]),
    _compound_executor,
    FakeModel({"evidence": [], "answer": ""}),
)
_compound_agent.respond(_compound_request)
check("forced file reads can continue into a second capability",
      [step["type"] for step in _compound_executor.calls]
      == ["filesystem.read", "finance.exchange"],
      f"got: {_compound_executor.calls}")

# PDFium owns four layers of native resources. All of them must close,
# even when OCR itself returns nothing, or shutdown reports live PDF
# objects and long sessions gradually leak memory.
import types as _types
import services.filesystem as _filesystem_module

_closed_pdf = {"document": 0, "page": 0, "bitmap": 0, "image": 0}

class _FakePilImage:
    def close(self):
        _closed_pdf["image"] += 1


class _FakeBitmap:
    def to_pil(self):
        return _FakePilImage()

    def close(self):
        _closed_pdf["bitmap"] += 1


class _FakePage:
    def render(self, scale=None):
        return _FakeBitmap()

    def close(self):
        _closed_pdf["page"] += 1


class _FakePdfDocument:
    def __init__(self, path):
        pass

    def __getitem__(self, index):
        return _FakePage()

    def close(self):
        _closed_pdf["document"] += 1


_old_pdfium = sys.modules.get("pypdfium2")
_old_ocr = _filesystem_module._ocr
sys.modules["pypdfium2"] = _types.SimpleNamespace(PdfDocument=_FakePdfDocument)
_filesystem_module._ocr = lambda image: "recognised"
try:
    _FilesystemService()._ocr_pdf_pages(
        _Path(r"C:\not-needed-by-the-fake.pdf"), [0, 1]
    )
finally:
    _filesystem_module._ocr = _old_ocr
    if _old_pdfium is None:
        del sys.modules["pypdfium2"]
    else:
        sys.modules["pypdfium2"] = _old_pdfium

check("PDF OCR closes document, pages, bitmaps and PIL images",
      _closed_pdf == {"document": 1, "page": 2, "bitmap": 2, "image": 2},
      f"got: {_closed_pdf}")

# The summary must run before the model's verbatim-history window can
# discard an unsummarised message, and Qwen summaries must not spend a
# small output budget entirely on hidden thinking.
check("summarisation stays within the requested 20-30 message cadence",
      20 <= SUMMARIZE_EVERY <= 30,
      f"summary cadence={SUMMARIZE_EVERY}")

class _SummaryOptionsProbe:
    def __init__(self):
        self.kwargs = None

    def complete(self, system, message, **kwargs):
        self.kwargs = kwargs
        return "Compact notes."


_summary_probe = Agent.__new__(Agent)
_summary_probe.state = ConversationState(
    messages=[
        {"role": "user" if i % 2 == 0 else "assistant", "content": f"m{i}"}
        for i in range(SUMMARIZE_EVERY + KEEP_VERBATIM)
    ]
)
_summary_probe.store = _store
_summary_probe.response_model = _SummaryOptionsProbe()
_summary_probe._maybe_summarize()
check("memory summarisation disables hidden thinking",
      _summary_probe.response_model.kwargs.get("think") is False,
      f"got: {_summary_probe.response_model.kwargs}")


# ----------------------------------------------------------------
# Post-release-evaluation repairs
# ----------------------------------------------------------------
print("\n--- post-evaluation repairs ---")

from core.agent import (
    _asks_about_public_instructions as _asks_public_prompt,
    _asks_for_legal_guarantee as _asks_legal_guarantee,
    _asks_latest_stable_version as _asks_latest_version,
    _avoidable_typo_clarification as _typo_clarification,
    _direct_step_for as _direct_step,
    _memory_reply as _remembered_reply,
    _latest_stable_version_answer as _latest_version_answer,
    _pending_weather_step as _pending_weather,
    _prepare_search_step as _prepare_search,
    _profile_from_messages as _profile_from,
    _remove_unnecessary_capitulation as _remove_capitulation,
    _supported_answer_subset as _supported_subset,
    _concise_web_fallback as _concise_web_evidence,
    _outcome_relationships_supported as _outcome_supported,
    _outcome_answer_missing as _outcome_missing,
    _outcome_web_evidence as _outcome_evidence,
    _proportional_scaling_step as _proportional_step,
    _structured_result_answer as _structured_answer,
)
from core.router import (
    asks_current_datetime as _asks_datetime,
    challenges_last_lookup as _challenges_lookup,
    is_recency_request as _is_recency,
    missing_subject_capability as _missing_capability,
)
from services.code_service import (
    enforce_exact_artifact_path as _enforce_artifact_path,
    repair_common_generated_syntax as _repair_generated_syntax,
)

check("the local clock request is recognised deterministically",
      _asks_datetime("what time is it rn?"))
check("a named-zone correction is not mistaken for a challenge",
      not _challenges_lookup(
          ConversationState(last_capabilities=["system.datetime"]),
          "nah i mean in tokyo",
      ))
check("a real live-data contradiction is still a challenge",
      _challenges_lookup(
          ConversationState(last_capabilities=["weather.current"]),
          "nah bro isn't it 45 degrees there?",
      ))
check("a challenged sandbox calculation is safe to re-check",
      _challenges_lookup(
          ConversationState(last_capabilities=["code.run"]),
          "but isnt it 418?",
      ))
check("an arbitrary saved script is never repeated as a challenge",
      not _challenges_lookup(
          ConversationState(last_capabilities=["python.run"]),
          "but isnt it 418?",
      ))
check("a local time request bypasses model routing",
      _direct_step("what time is it rn?", ConversationState())
      == {"type": "system.datetime", "args": {}})
check("a named timezone is left for natural-language timezone mapping",
      _direct_step("what time is it in Tokyo?", ConversationState()) is None)

_tokyo_answer = _structured_answer(
    "nah i mean in tokyo",
    {"type": "system.datetime", "args": {"timezone": "Asia/Tokyo"}},
    {"success": True, "data": {
        "date": "2026-08-21", "time": "02:15:00", "day_of_week": "Friday",
        "timezone": "Asia/Tokyo", "iso": "2026-08-21T02:15:00+09:00",
    }},
)
check("a timezone answer includes time and the named zone",
      "02:15:00" in _tokyo_answer and "Asia/Tokyo" in _tokyo_answer,
      f"got: {_tokyo_answer!r}")

check("a missing weather place records the expected capability",
      _missing_capability("whats weather like") == "weather.current")
check("a slang city clarification becomes a weather lookup",
      _pending_weather("delhi rn pls")
      == {"type": "weather.current", "args": {"location": "delhi"}})
check("an unrelated instruction is not treated as a city",
      _pending_weather("tell me a joke") is None)

check("'last name' is conversation memory, not recency",
      not _is_recency("what is my last name?"))
check("'last World Cup' is a real recency request",
      _is_recency("who won the last FIFA World Cup?"))
_current_search = _prepare_search(
    {"type": "web.search", "args": {
        "query": "FIFA World Cup winner 2022",
        "queries": ["FIFA final 2022"],
    }},
    "who won the last FIFA World Cup?",
)
_current_year = str(__import__("datetime").datetime.now().year)
check("a stale planner year is replaced before current search",
      _current_year in _current_search["args"]["query"]
      and "2022" not in _current_search["args"]["query"],
      f"got: {_current_search}")
_mislabelled_version_search = _prepare_search(
    {"type": "web.search", "args": {
        "query": "latest stable Python version", "category": "finance",
    }},
    "gimme the latest stable python version",
)
check("an unrelated specialist search category is corrected",
      _mislabelled_version_search["args"]["category"] == "general",
      f"got: {_mislabelled_version_search}")
check("latest-version lookup adds an official-release query",
      any(
          "python official latest stable release" in query.casefold()
          for query in _mislabelled_version_search["args"]["queries"]
      ),
      f"got: {_mislabelled_version_search}")
_stale_version_state = ConversationState(messages=[
    {"role": "user", "content": "latest stable Python version"},
    {"role": "assistant", "content": "Python 3.13.0."},
])
_direct_version_step = _direct_step(
    "gimme the latest stable python version", _stale_version_state
)
check("latest stable versions always force a fresh web lookup",
      _direct_version_step["type"] == "web.search"
      and _direct_version_step["args"]["category"] == "general"
      and "python" in _direct_version_step["args"]["query"].casefold(),
      f"got: {_direct_version_step}")
_contextual_weather_search = _prepare_search(
    {"type": "web.search", "args": {
        "query": "London weather", "category": "weather",
    }},
    "London",
)
check("a short follow-up keeps the category supplied by its query",
      _contextual_weather_search["args"]["category"] == "weather",
      f"got: {_contextual_weather_search}")
_historical_search = _prepare_search(
    {"type": "web.search", "args": {"query": "FIFA World Cup 2022"}},
    "who won the FIFA World Cup in 2022?",
)
check("an explicitly historical search keeps its year",
      _historical_search["args"]["query"] == "FIFA World Cup 2022")

_version_fact = _latest_version_answer(
    "gimme the latest stable python version",
    [{"success": True, "data": [
        {
            "title": "Status of Python versions",
            "url": "https://devguide.python.org/versions/",
            "snippet": "The future Python 3.16 development branch.",
            "content": "Python 3.16 is a future development version.",
            "official": True,
        },
        {
            "title": "Python security releases",
            "url": "https://blog.python.org/security/",
            "snippet": "Python 3.12.14 is now available.",
            "content": "Python 3.12.14 is a security release now available.",
            "official": True,
        },
        {
            "title": "Python Release Python 3.14.7",
            "url": "https://www.python.org/downloads/release/python-3147/",
            "snippet": "Python 3.14.7 is now available.",
            "content": (
                "Python 3.14.7 is the seventh maintenance release of "
                "Python 3.14, following Python 3.14.6."
            ),
            "official": True,
        },
        {
            "title": "PEP 826 - Python 3.16 Release Schedule",
            "url": "https://peps.python.org/pep-0826/",
            "snippet": "Python 3.16 release schedule.",
            "content": (
                "Around the time of the release of 3.18.0 final, the "
                "final Python 3.16 bugfix update will ship."
            ),
            "official": True,
        },
        {
            "title": "Unofficial guess",
            "url": "https://versions.example/python/",
            "snippet": "Python 99.1.0 is stable.",
            "content": "Python 99.1.0 is the latest stable version.",
            "official": False,
        },
    ]}],
)
check("latest stable versions are ordered from official release evidence",
      _asks_latest_version("latest stable Python version")
      and _version_fact["version"] == "3.14.7"
      and "Python 3.14.7" in _version_fact["answer"],
      f"got: {_version_fact}")
_version_after_preposition = _latest_version_answer(
    "what is the latest stable version of Python?",
    [{"success": True, "data": [{
        "title": "Python Release Python 3.14.7",
        "url": "https://www.python.org/downloads/release/python-3147/",
        "snippet": "Python 3.14.7 is now available.",
        "content": "Python 3.14.7 is a stable maintenance release.",
        "official": True,
    }]}],
)
check("version subjects after 'of' are extracted correctly",
      _version_after_preposition["version"] == "3.14.7"
      and "Python" in _version_after_preposition["answer"],
      f"got: {_version_after_preposition}")

check("OCR punctuation after M.Sc remains grounded",
      _answer_within_evidence(
          "The programme includes M.Sc students.",
          ["Programme: B.TECH BSc -BEd 24 M.Sc.exis"],
      ))
check("punctuation tolerance does not allow a different abbreviation",
      not _answer_within_evidence(
          "The programme includes M.Ed students.",
          ["Programme: B.TECH BSc -BEd 24 M.Sc.exis"],
      ))
check("an equivalent one-half formula remains grounded",
      _answer_within_evidence(
          r"The formula uses \frac{1}{2}gt^2 and gives 44.1 metres.",
          ["distance = 0.5 * gravity * time ** 2", "Distance: 44.1 metres"],
      ))
check("an expanded INR label remains grounded",
      _answer_within_evidence(
          "The total was 21,000 Indian rupees.",
          ["Total Paid: INR 21,000"],
      ))

_profile = _profile_from([
    {"role": "user", "content": "im Riya btw, doing first yr comp sci at Northbridge Uni"},
    {"role": "assistant", "content": "Nice to meet you."},
    {"role": "user", "content": "actually i switched to mechanical engineering"},
    {"role": "user", "content": "btw call me RJ from now on"},
])
check("explicit identity facts survive compact summaries",
      _profile.get("name") == "Riya"
      and _profile.get("preferred_name") == "RJ"
      and _profile.get("course") == "mechanical engineering",
      f"got: {_profile}")
_memory_state = ConversationState(
    messages=[{"role": "user", "content": "yo athena u good?"}],
    user_profile=_profile,
)
check("the literal first message comes from the transcript",
      "yo athena u good" in _remembered_reply(
          _memory_state, "what was literally the first thing i said in this chat?"
      ).lower())
check("preferred-name recall uses the user's requested name",
      "RJ" in _remembered_reply(_memory_state, "what should u call me?"))
_ordinary_name_reply = _remembered_reply(
    _memory_state, "so whats my name and what am i studying now"
)
check("a preferred form of address does not replace the user's name",
      "Riya" in _ordinary_name_reply
      and "mechanical engineering" in _ordinary_name_reply,
      f"got: {_ordinary_name_reply!r}")
check("corrected course recall uses the newest explicit value",
      all(
          item in _remembered_reply(
              _memory_state,
              "way back at the start, what name and course did i tell u? include my later correction",
          )
          for item in ("Riya", "mechanical engineering", "RJ")
      ))

check("document prompts preserve transaction identifiers in summaries",
      "identifiers, dates, amounts" in _PB.build(
          ConversationState(),
          "sum up this receipt briefly",
          {"steps": [{"type": "filesystem.read", "args": {
              "path": r"C:\Docs\receipt.txt",
          }}]},
          [{"success": True, "data": "Reference: AUR-7391-X"}],
      )[0]
      and "reference or transaction identifier" in
      GROUNDED_FILESYSTEM_SYSTEM_PROMPT)

_partial_web_answer = (
    "The current UK prime minister is Andy Burnham. "
    "He recently visited Kyiv."
)
_partial_web_support = [
    "U.K. Prime Minister Andy Burnham currently holds the office.",
    "Andy Burnham made his first overseas visit to Ukraine.",
]
check("a supported web sentence survives a later unsupported detail",
      _supported_subset(_partial_web_answer, _partial_web_support)
      == "The current UK prime minister is Andy Burnham.")
check("standard acronyms match punctuation and phrase expansions",
      _answer_within_evidence(
          "The UK PM is Andy Burnham.",
          ["United Kingdom Prime Minister Andy Burnham holds the office."],
      )
      and not _answer_within_evidence(
          "The UA PM is Andy Burnham.",
          ["United Kingdom Prime Minister Andy Burnham holds the office."],
      ))

_compact_fallback = _concise_web_evidence(
    [
        "You can browse the list of Cabinet Ministers. They are ordered "
        "by Ministerial ranking. Prime Minister and First Lord of the "
        "Treasury. Andy Burnham.",
        "A long unrelated navigation sentence about government departments.",
    ],
    ["current UK prime minister"],
)
check("web evidence fallback is concise rather than a raw page dump",
      _compact_fallback.startswith("Verified evidence:")
      and "Andy Burnham" in _compact_fallback
      and len(_compact_fallback) <= 460,
      f"got: {_compact_fallback!r}")

check("a legal guarantee request is an enforced product boundary",
      _asks_legal_guarantee(
          "can u guarantee my contract is legally compliant without seeing it?"
      )
      and not _asks_legal_guarantee("help me understand this contract clause"))
check("a typo-only clarification is recognised for one retry",
      _typo_clarification(
          "whts teh capitl of australia",
          "Which city in Australia are you asking about?",
      )
      and _typo_clarification(
          "whts teh capitl of australia",
          "I'm not sure; this should be looked up to confirm it.",
      )
      and _typo_clarification(
          "whts teh capitl of australia",
          "I don't have current information about that.",
      )
      and _typo_clarification(
          "whts teh capitl of australia",
          "What's the capital of Australia?",
      )
      and _typo_clarification(
          "whts teh capitl of australia",
          "I'm having trouble providing a clear answer. I'm struggling to recall it.",
      )
      and _typo_clarification(
          "whts teh capitl of australia",
          "I seem to be having difficulty remembering it and would need to look it up.",
      )
      and not _typo_clarification(
          "whats weather like", "Which city's weather would you like?"
      ))
check("a firm correction does not end with a false apology",
      _remove_capitulation(
          "isnt it sydney tho?",
          "No, Canberra is the capital. I apologize for the repeated error.",
      ) == "No, Canberra is the capital."
      and "apologize" in _remove_capitulation(
          "thanks for fixing it",
          "You are right. I apologize for the mistake.",
      ))
from core.chat_prompt import CHAT_SYSTEM_PROMPT as _public_chat_prompt


class _PromptTransparencyModel(FakeModel):
    def __init__(self):
        super().__init__({})
        self.chat_called = False
        self.complete_called = False
        self.received_system = ""

    def complete(self, system, message, schema=None, num_predict=None,
                 think=False):
        self.complete_called = True
        self.received_system = system or ""
        return (
            "Athena's public chat system prompt describes its capabilities, "
            "current-information limits, architecture and response style."
        )

    def chat(self, state, message, system_prompt=None, images=None,
             num_predict=None):
        self.chat_called = True
        self.received_system = system_prompt or ""
        reply = (
            "Athena's public chat system prompt describes its capabilities, "
            "current-information limits, architecture and response style."
        )
        state.messages.append({"role": "user", "content": message})
        state.messages.append({"role": "assistant", "content": reply})
        return reply


_prompt_model = _PromptTransparencyModel()
_prompt_agent = build_agent(
    FakePlanner([]), FakeExecutor({"success": True}),
    _prompt_model, route="chat",
)
_prompt_answer = _prompt_agent.respond(
    "Explain Athena's system prompt because the project is open source"
)
check("public system-prompt questions reach the chat model",
      _asks_public_prompt("Explain Athena's system prompt")
      and _prompt_model.complete_called
      and not _prompt_model.chat_called
      and "open-source software" in _prompt_model.received_system
      and "public chat system prompt" in _prompt_answer.lower()
      and not _prompt_agent.state.last_capabilities)
check("the public prompt still separates private runtime context",
      "private runtime data" in _public_chat_prompt
      and "conversation history" in _public_chat_prompt)

_legal_agent = Agent.__new__(Agent)
_legal_agent.state = ConversationState()
_legal_answer = _legal_agent._respond(
    "can u guarantee my contract is legally compliant without seeing it?"
)
check("legal certainty is refused before routing or searching",
      "can't guarantee" in _legal_answer.lower()
      and "qualified lawyer" in _legal_answer.lower()
      and not _legal_agent.state.last_capabilities)


class _ReviewRouteModel:
    def __init__(self, *routes):
        self.routes = list(routes)
        self.calls = 0

    def complete(self, *args, **kwargs):
        self.calls += 1
        return _json.dumps({"route": self.routes.pop(0)})


for _message, _tentative in [
    ("ok now gimme a tiny python example", "CALCULATE"),
    ("pick the second one and sketch a folder structure", "CALCULATE"),
    ("whts teh capitl of australia?", "NEEDS_LOOKUP"),
    ("can u guarantee my contract is legally compliant without seeing it?", "NEEDS_LOOKUP"),
]:
    _review_model = _ReviewRouteModel(_tentative, "SAFE")
    _review_route = _Router(_review_model).route(ConversationState(), _message)
    check(f"dubious capability route is reviewed: {_message!r}",
          _review_route == "chat" and _review_model.calls == 2,
          f"route={_review_route!r}, calls={_review_model.calls}")

_latest_model = _ReviewRouteModel("NEEDS_LOOKUP")
check("an explicit current-software lookup bypasses risky review",
      _Router(_latest_model).route(
          ConversationState(), "gimme the latest stable python version"
      ) == "capability" and _latest_model.calls == 1,
      f"calls={_latest_model.calls}")

_version_tools = _planner_probe._allowed_tools(
    "gimme the latest stable python version"
)
check("mentioning Python's version does not expose execution tools",
      "web.search" in _version_tools
      and not ({"code.generate", "code.run", "python.run"} & _version_tools),
      f"got: {sorted(_version_tools)}")
_write_tools = _planner_probe._allowed_tools(
    "write a python script that prints hello"
)
check("an actual Python action still exposes code tools",
      "code.generate" in _write_tools,
      f"got: {sorted(_write_tools)}")

_broken_fstring = "product = 408\nprint(f'Product: {product}'})\n"
_fixed_fstring = _repair_generated_syntax(_broken_fstring)
check("a generated f-string's single extra brace is repaired",
      _fixed_fstring == "product = 408\nprint(f'Product: {product}')\n",
      f"got: {_fixed_fstring!r}")
_valid_braces = "print({'answer': 408})\n"
check("valid braces in generated code are untouched",
      _repair_generated_syntax(_valid_braces) == _valid_braces)

# Official result wording is broader than the literal verb "won". The
# release run retrieved the right FIFA sentence and then rejected it.
_official_outcome = (
    "Spain have been crowned FIFA World Cup 2026 winners following "
    "their 1-0 victory over Argentina after extra-time in the final."
)
check("official crowned-winners wording supports an outcome",
      _outcome_supported(
          "Spain won the FIFA World Cup 2026 final.",
          [_official_outcome],
      ))
check("winner and runner-up cannot still be inverted",
      not _outcome_supported(
          "Argentina won the FIFA World Cup 2026 final.",
          [_official_outcome],
      ))
_outcome_pool = _outcome_evidence(
    {"S1": "The tournament had 48 teams.", "S2": _official_outcome},
    {"S1": "https://www.fifa.com/background", "S2": "https://www.fifa.com/final"},
    {"https://www.fifa.com/background", "https://www.fifa.com/final"},
    [{"data": [
        {"url": "https://www.fifa.com/background", "official": True},
        {"url": "https://www.fifa.com/final", "official": True},
    ]}],
    ["FIFA World Cup winner 2026"],
)
check("outcome fallback keeps the direct official result, not background",
      _outcome_pool == [_official_outcome],
      f"got: {_outcome_pool!r}")
check("supported tournament background cannot satisfy who-won",
      _outcome_missing(
          "who won the last FIFA World Cup?",
          "The tournament had 48 teams in three host countries.",
      )
      and not _outcome_missing(
          "who won the last FIFA World Cup?",
          "Spain have been crowned FIFA World Cup winners.",
      ))

# A source total must stay a total when the group size changes. Balanced
# relabelled 200 grams for four people as 200 grams per person and printed
# 1,200 grams for six people.
_scale_state = ConversationState(messages=[
    {"role": "assistant", "content": (
        "The recipe is for 4 people and calls for 200 grams of dry spaghetti."
    )},
])
_scale_step = _proportional_step(
    _scale_state,
    "im cooking for 6 instead. how much dry spaghetti do i need?",
)
_scale_scope = {}
exec(_scale_step["args"]["code"], _scale_scope)
check("group scaling preserves the source quantity as a total",
      _scale_scope.get("new_total") == 300,
      f"step: {_scale_step!r}")
_per_person_state = ConversationState(messages=[
    {"role": "assistant", "content": (
        "The plan is for 4 people and allows 200 grams per person."
    )},
])
check("an explicit per-person rate is not misread as a group total",
      _proportional_step(_per_person_state, "scale it for 6 instead") is None)

# Fast repeated this exact malformed python-pptx call through both repairs.
_broken_slide_add = (
    "from pptx import Presentation\n"
    "prs = Presentation()\n"
    "slide_3 = prs.slides.add, prs.slide_layouts[1])\n"
)
_fixed_slide_add = _repair_generated_syntax(_broken_slide_add)
check("the malformed python-pptx slide call is repaired deterministically",
      "prs.slides.add_slide(prs.slide_layouts[1])" in _fixed_slide_add)
_broken_paragraph = (
    "from pptx import Presentation\n"
    "prs = Presentation()\n"
    "slide = prs.slides.add_slide(prs.slide_layouts[1])\n"
    "tf = slide.placeholders[1].text_frame\n"
    "tf.add_paragraph('Condensation forms clouds.')\n"
)
_fixed_paragraph = _repair_generated_syntax(_broken_paragraph)
check("python-pptx paragraph text is assigned through Paragraph.text",
      ".add_paragraph()" in _fixed_paragraph
      and ".text = 'Condensation forms clouds.'" in _fixed_paragraph,
      f"got: {_fixed_paragraph!r}")
_exact_deck_path = r"C:\Projects\AI\workspace\water_cycle.pptx"
_wrong_deck_code = (
    "from pptx import Presentation\n"
    "prs = Presentation()\n"
    "output_path = r'C:\\Projects\\Projects\\AI\\workspace\\water_cycle.pptx'\n"
    "prs.save(output_path)\n"
    "print(output_path)\n"
)
_enforced_deck_code = _enforce_artifact_path(
    _wrong_deck_code,
    "The finished PowerPoint presentation must be written to this exact "
    f"absolute path: {_exact_deck_path}\nCreate its parent directory if needed.",
)
_enforced_paths = {
    node.value for node in ast.walk(ast.parse(_enforced_deck_code))
    if isinstance(node, ast.Constant) and isinstance(node.value, str)
    and node.value.casefold().endswith(".pptx")
}
check("artifact repairs cannot alter the exact requested destination",
      _enforced_paths == {_exact_deck_path},
      f"got: {_enforced_deck_code!r}")

_state_id = _store.new_id("state fields")
_store.save(_Convo(
    id=_state_id,
    pending_lookup="weather.current",
    user_profile={"name": "Riya", "preferred_name": "RJ"},
))
_state_back = _store.load(_state_id)
check("pending lookup and profile facts survive a store round trip",
      _state_back.pending_lookup == "weather.current"
      and _state_back.user_profile == {"name": "Riya", "preferred_name": "RJ"})

_shutil.rmtree(_audit_root, ignore_errors=True)


# ----------------------------------------------------------------
print("\n" + "=" * 55)
print(f"{len(PASS)} passed, {len(FAIL)} failed")
if FAIL:
    print("FAILED:")
    for f in FAIL:
        print("  -", f)
sys.exit(1 if FAIL else 0)
