"""
Filesystem service.

All filesystem access goes through this service.
"""

import io
import shutil
from pathlib import Path

from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import xlrd
import pytesseract
from PIL import Image
import os
import re
from datetime import datetime, timezone
from config import MAX_FILE_SIZE as CONFIG_MAX_FILE_SIZE

def _find_tesseract() -> str | None:
    """Locate the OCR binary without assuming where it was installed.

    Checked in order of how deliberate each source is: an explicit
    environment variable beats a copy on PATH, which beats a guess at
    the usual install directories.

    Returns None when there is nothing to find. OCR is an extra attempt
    on top of normal text extraction, so its absence should cost a
    reader nothing - the alternative, a hardcoded path to one machine,
    meant every other machine raised on the first scanned page.
    """

    explicit = os.environ.get("TESSERACT_CMD")

    if explicit and Path(explicit).exists():
        return explicit

    found = shutil.which("tesseract")

    if found:
        return found

    candidates = [
        # Windows, per-user and machine-wide installs.
        Path(os.environ.get("LOCALAPPDATA", "")) / "Programs/Tesseract-OCR/tesseract.exe",
        Path(os.environ.get("PROGRAMFILES", "")) / "Tesseract-OCR/tesseract.exe",
        Path(os.environ.get("PROGRAMFILES(X86)", "")) / "Tesseract-OCR/tesseract.exe",
        # macOS (Homebrew, Intel and Apple silicon) and Linux.
        Path("/opt/homebrew/bin/tesseract"),
        Path("/usr/local/bin/tesseract"),
        Path("/usr/bin/tesseract"),
    ]

    for candidate in candidates:
        try:
            if candidate.is_file():
                return str(candidate)
        except OSError:
            # An unreadable or malformed path is simply not a match.
            continue

    return None


_TESSERACT = _find_tesseract()

if _TESSERACT:
    pytesseract.pytesseract.tesseract_cmd = _TESSERACT
else:
    print(
        "[OCR] Tesseract not found - scanned pages and images inside "
        "documents will not be read. Install it, or set TESSERACT_CMD "
        "to its full path."
    )

def _ocr(image) -> str:
    """OCR a PIL image, returning "" rather than raising.

    OCR is always an extra attempt on top of normal extraction, so a
    failure here should cost the caller nothing.
    """

    try:
        return (pytesseract.image_to_string(image) or "").strip()

    except Exception as error:
        print(f"[OCR] failed: {error}")
        return ""


def _ocr_bytes(blob: bytes) -> str:
    """OCR raw image bytes pulled out of a document."""

    try:
        with Image.open(io.BytesIO(blob)) as image:
            return _ocr(image)

    except Exception:
        # Not a format PIL can open, or not really an image.
        return ""


class FilesystemService:

    # A page yielding less than this is treated as having no real text,
    # and is put through OCR. Not zero: a scanned page often still
    # carries a stray header or page number in a text layer.
    OCR_MIN_TEXT_PER_PAGE = 40

    # Rendering resolution for OCR. 300 is the usual floor for reliable
    # character recognition; lower loses small print on receipts.
    OCR_RENDER_DPI = 300

    # Rendering is the slow part - roughly a second a page - so a long
    # scanned document is capped rather than left to run for minutes.
    OCR_MAX_PAGES = 20

    # Folders to skip during search - noisy, huge, or irrelevant
    SKIP_DIRS = {
        "AppData", ".git", "node_modules", "__pycache__",
        ".venv", "venv", "$RECYCLE.BIN", "System Volume Information",
        ".cache"
    }

    # Secrets are not ordinary documents. Athena searches a user's home
    # directory, so reading one into a model prompt by accident would be a
    # much worse failure than refusing a legitimate edge case. The public
    # .env.example remains readable because it contains placeholders only.
    SENSITIVE_NAMES = {
        ".env", ".npmrc", ".pypirc", ".netrc",
        "id_rsa", "id_ed25519", "credentials.json", "secrets.json",
    }
    SENSITIVE_SUFFIXES = {".pem", ".key", ".p12", ".pfx", ".jks"}

    @classmethod
    def _is_sensitive(cls, file: Path) -> bool:
        name = file.name.lower()

        if name in cls.SENSITIVE_NAMES:
            return True

        if name.startswith(".env.") and name != ".env.example":
            return True

        return file.suffix.lower() in cls.SENSITIVE_SUFFIXES

    def _normalize(self, text: str) -> str:
        return re.sub(r"[\s_-]+", "", text.lower())

    def search(self, name: str, max_results: int = 10):

        try:

            if not name or not name.strip():

                return {
                    "success": False,
                    "error": "What is the file's name?"
                }

            query = self._normalize(name)
            home = Path.home()

            matches = []

            for root, dirs, files in os.walk(home):

                dirs[:] = [d for d in dirs if d not in self.SKIP_DIRS]

                for filename in files:

                    if self._is_sensitive(Path(filename)):
                        continue

                    if query in self._normalize(filename):

                        matches.append(Path(root) / filename)

                        if len(matches) >= max_results:
                            break

                if len(matches) >= max_results:
                    break

            if not matches:

                return {
                    "success": False,
                    "error": f"No file found matching '{name}'."
                }

            if len(matches) > 1:

                paths = "\n".join(str(m) for m in matches)

                # "matches" is returned in the same order it is shown to
                # the user, so a follow-up like "the first one" can be
                # resolved against this list directly rather than having
                # the planner re-derive a path from conversation text.
                return {
                    "success": False,
                    "error": (
                        f"Found multiple files matching '{name}':\n{paths}\n\n"
                        f"Which one did you mean?"
                    ),
                    "matches": [str(m) for m in matches]
                }

            result = self.read(str(matches[0]))

            if result.get("success"):
                result["resolved_path"] = str(matches[0])

            return result

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }

    # Maximum file size, configurable in .env for unusually large documents.
    MAX_FILE_SIZE = CONFIG_MAX_FILE_SIZE

    # Plain text and code/config files
    TEXT_EXTENSIONS = {
        ".txt", ".md", ".rst", ".log", ".csv", ".tsv",
        ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp",
        ".h", ".hpp", ".go", ".rs", ".rb", ".php", ".swift", ".kt",
        ".sh", ".bat", ".ps1", ".sql",
        ".html", ".htm", ".css", ".scss",
        ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
        ".xml", ".env"
    }

    # Real documents requiring dedicated parsing
    DOCUMENT_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".xlsx",
        ".pptx",
        ".xls",
    }

    IMAGE_EXTENSIONS = {
        ".jpg", ".jpeg", ".png"
    }

    # Extensions we know are binary/non-text and have no reader for -
    # skip these outright rather than trying to decode garbage.
    KNOWN_BINARY_EXTENSIONS = {
        ".exe", ".dll", ".so", ".bin", ".dat",
        ".gif", ".bmp", ".ico", ".webp",
        ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv",
        ".zip", ".rar", ".7z", ".tar", ".gz",
        ".sqlite", ".db",
        ".pyc", ".class",
        ".ttf", ".otf", ".woff", ".woff2",
        ".doc",  # legacy Word format - not supported
    }

    def exists(self, path: str):

        try:

            return {
                "success": True,
                "data": Path(path).exists()
            }

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }

    def read(self, path: str):

        try:

            file = Path(path)

            if not file.exists():

                return {
                    "success": False,
                    "error": "File does not exist."
                }

            if not file.is_file():

                return {
                    "success": False,
                    "error": "Path is not a file."
                }

            if self._is_sensitive(file):

                return {
                    "success": False,
                    "error": (
                        f"'{file.name}' looks like a credentials or secrets file, "
                        "so Athena will not read it."
                    )
                }

            suffix = file.suffix.lower()

            if suffix in self.KNOWN_BINARY_EXTENSIONS:

                return {
                    "success": False,
                    "error": f"'{suffix}' is a binary file type and can't be read as text."
                }

            size = file.stat().st_size

            if size > self.MAX_FILE_SIZE:

                return {
                    "success": False,
                    "error": f"File is too large ({size} bytes)."
                }

            if suffix == ".pdf":
                text = self._read_pdf(file)

            elif suffix == ".docx":
                text = self._read_docx(file)

            elif suffix == ".xlsx":
                text = self._read_xlsx(file)

            elif suffix == ".xls":
                text = self._read_xls(file)

            elif suffix == ".pptx":
                text = self._read_pptx(file)

            elif suffix in self.IMAGE_EXTENSIONS:
                text = self._read_image(file)

            else:

                raw = file.read_bytes()

                try:
                    text = raw.decode("utf-8")
                except UnicodeDecodeError:
                    text = raw.decode("utf-8", errors="replace")
                    replaced_ratio = text.count("\ufffd") / max(len(text), 1)
                    if replaced_ratio > 0.05:
                        return {
                            "success": False,
                            "error": f"'{file.name}' doesn't appear to be a readable text file."
                        }

            # A PDF of scanned pages parses fine and yields no text at
            # all. Passed on as an empty success, that reached the user
            # as "I couldn't find that", which sounds like the answer
            # wasn't in the document rather than like nothing could be
            # read from it.
            if suffix in self.DOCUMENT_EXTENSIONS and not text.strip():

                return {
                    "success": False,
                    "error": (
                        f"'{file.name}' opened, but no text could be read from it - "
                        f"not from the document itself and not from any image in it, "
                        f"which were put through text recognition as well. It may be "
                        f"blank, or a picture with no legible writing."
                    )
                }

            return {
                "success": True,
                "data": text
            }

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }

    def _read_pdf(self, file: Path) -> str:
        """Text layer first, OCR for whatever it doesn't cover.

        Plenty of real documents have no text to extract. A scan is
        images of pages; an exported receipt can have its text
        converted to vector outlines, which draws identically and
        leaves nothing to read. Both used to come back empty. Pages
        that yield no text are rendered and read with OCR instead.
        """

        reader = PdfReader(str(file))

        pages = []
        needs_ocr = []

        for index, page in enumerate(reader.pages):

            text = (page.extract_text() or "").strip()
            pages.append(text)

            if len(text) < self.OCR_MIN_TEXT_PER_PAGE:
                needs_ocr.append(index)

        if needs_ocr:

            print(f"[OCR] {len(needs_ocr)} page(s) of '{file.name}' have no text - reading them as images")

            for index, text in self._ocr_pdf_pages(file, needs_ocr).items():

                if not text:
                    continue

                pages[index] = f"{pages[index]}\n{text}".strip() if pages[index] else text

        return "\n\n".join(p for p in pages if p).strip()

    def _ocr_pdf_pages(self, file: Path, indexes) -> dict:
        """Render the given pages and OCR them.

        Rendering needs pypdfium2. Embedded images are tried first
        through pypdf, which needs no renderer and is much faster - it
        covers the common scan, where the page is one big image.
        """

        results = {}
        remaining = []

        try:
            reader = PdfReader(str(file))

            for index in indexes:

                text = ""

                try:
                    for image in reader.pages[index].images:
                        text = f"{text}\n{_ocr_bytes(image.data)}".strip()

                except Exception:
                    text = ""

                if len(text) >= self.OCR_MIN_TEXT_PER_PAGE:
                    results[index] = text
                else:
                    remaining.append(index)

        except Exception as error:
            print(f"[OCR] embedded-image pass failed: {error}")
            remaining = list(indexes)

        if not remaining:
            return results

        # Anything left is drawn rather than embedded, so the page has
        # to be rendered before it can be read.
        try:
            import pypdfium2

        except ImportError:
            print("[OCR] pypdfium2 is not installed, so drawn pages can't be rendered")
            return results

        document = None

        try:
            document = pypdfium2.PdfDocument(str(file))

            for index in remaining[:self.OCR_MAX_PAGES]:
                page = None
                bitmap = None
                image = None

                try:
                    page = document[index]
                    bitmap = page.render(scale=self.OCR_RENDER_DPI / 72)
                    image = bitmap.to_pil()
                    results[index] = _ocr(image)

                except Exception as error:
                    # One malformed page should not prevent OCR on the
                    # remaining pages of an otherwise readable document.
                    print(f"[OCR] page {index + 1} failed: {error}")

                finally:
                    if image is not None:
                        image.close()
                    if bitmap is not None:
                        bitmap.close()
                    if page is not None:
                        page.close()

            if len(remaining) > self.OCR_MAX_PAGES:
                print(f"[OCR] stopped after {self.OCR_MAX_PAGES} pages")

        except Exception as error:
            print(f"[OCR] page rendering failed: {error}")

        finally:
            if document is not None:
                document.close()

        return results

    @staticmethod
    def _ocr_embedded_parts(package) -> str:
        """OCR every image stored inside an Office file.

        A receipt or screenshot pasted into a document is invisible to
        normal text extraction, so anything written in it would be
        missed even though the user can plainly see it.
        """

        found = []

        try:
            for part in package.part.package.iter_parts():

                if "image" not in getattr(part, "content_type", ""):
                    continue

                text = _ocr_bytes(part.blob)

                if text:
                    found.append(text)

        except Exception as error:
            print(f"[OCR] embedded images could not be read: {error}")

        if not found:
            return ""

        return "\n\n".join(["[Text found in embedded images:]"] + found)

    def _read_docx(self, file: Path) -> str:

        doc = Document(str(file))

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    paragraphs.append("\t".join(cells))

        text = "\n".join(paragraphs)

        images = self._ocr_embedded_parts(doc)

        return f"{text}\n\n{images}".strip() if images else text

    def _read_xlsx(self, file: Path) -> str:

        workbook = load_workbook(str(file), data_only=True)

        sections = []

        for sheet in workbook.worksheets:

            rows = []

            for row in sheet.iter_rows(values_only=True):

                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))

            sections.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows))

            # openpyxl keeps sheet images on a private attribute; there
            # is no public accessor, so this is guarded rather than
            # relied upon.
            for image in getattr(sheet, "_images", []) or []:

                try:
                    blob = image.ref.getvalue() if hasattr(image.ref, "getvalue") else image.ref

                    text = _ocr_bytes(blob) if isinstance(blob, bytes) else _ocr(blob)

                    if text:
                        sections.append(f"[Text found in an image on {sheet.title}:]\n{text}")

                except Exception:
                    continue

        return "\n\n".join(sections)

    def _read_xls(self, file: Path) -> str:

        workbook = xlrd.open_workbook(str(file))

        sections = []

        for sheet in workbook.sheets():

            rows = []

            for row_idx in range(sheet.nrows):

                cells = [str(c) for c in sheet.row_values(row_idx)]
                rows.append("\t".join(cells))

            sections.append(f"--- Sheet: {sheet.name} ---\n" + "\n".join(rows))

        return "\n\n".join(sections)

    def _read_pptx(self, file: Path) -> str:

        prs = Presentation(str(file))

        slides_text = []

        for i, slide in enumerate(prs.slides, start=1):

            parts = [f"--- Slide {i} ---"]

            for shape in slide.shapes:

                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs)
                        if line.strip():
                            parts.append(line)

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells]
                        parts.append("\t".join(cells))

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                parts.append(f"[Notes: {slide.notes_slide.notes_text_frame.text}]")

            slides_text.append("\n".join(parts))

        text = "\n\n".join(slides_text)

        images = self._ocr_embedded_parts(prs)

        return f"{text}\n\n{images}".strip() if images else text

    def _read_image(self, file: Path) -> str:

        image = Image.open(file)

        text = pytesseract.image_to_string(image)

        if not text.strip():
            return "[No readable text found in this image via OCR.]"

        return text

    # write(), list(), exists(), info() stay exactly as they are

    def write(self, path: str, content: str, overwrite: bool = False):

        try:

            if not path or not path.strip():

                return {
                    "success": False,
                    "error": "No filename was given. What would you like to name the file?"
                }

            file = Path(path)

            if not file.suffix:
                file = file.with_suffix(".py")

            path = str(file)

            if file.exists() and not overwrite:

                return {
                    "success": False,
                    "error": (
                        f"A file already exists at '{path}'. "
                        f"Would you like me to overwrite it, or save it "
                        f"under a different filename?"
                    ),
                    "conflict": True
                }

            if file.suffix.lower() not in self.TEXT_EXTENSIONS:

                return {
                    "success": False,
                    "error": f"Unsupported file type: {file.suffix}"
                }

            size = len(content.encode("utf-8"))

            if size > self.MAX_FILE_SIZE:

                return {
                    "success": False,
                    "error": f"Content is too large ({size} bytes)."
                }

            file.parent.mkdir(parents=True, exist_ok=True)

            file.write_text(content, encoding="utf-8")

            return {
                "success": True,
                "data": {
                    "path": str(file.resolve()),
                    "bytes_written": size
                }
            }

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }

    def list(self, path: str):

        try:

            directory = Path(path)

            if not directory.exists():

                return {
                    "success": False,
                    "error": "Directory does not exist."
                }

            if not directory.is_dir():

                return {
                    "success": False,
                    "error": "Path is not a directory."
                }

            items = []

            for item in sorted(directory.iterdir()):

                items.append(
                    {
                        "name": item.name,
                        "type": (
                            "directory"
                            if item.is_dir()
                            else "file"
                        )
                    }
                )

            return {
                "success": True,
                "data": items
            }

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }

    def info(self, path: str):

        try:

            file = Path(path)

            if not file.exists():

                return {
                    "success": False,
                    "error": "Path does not exist."
                }

            stat = file.stat()

            return {
                "success": True,
                "data": {
                    "name": file.name,
                    "path": str(file.resolve()),
                    "size": stat.st_size,
                    "extension": file.suffix,
                    "is_file": file.is_file(),
                    "is_directory": file.is_dir(),
                    "modified_at": datetime.fromtimestamp(
                        stat.st_mtime, tz=timezone.utc
                    ).astimezone().isoformat(timespec="seconds"),
                }
            }

        except Exception as e:

            return {
                "success": False,
                "error": str(e)
            }
